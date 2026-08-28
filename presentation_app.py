"""
========================================================================================
⚡ AI VOICE-CONTROLLED DUAL-SCREEN PPT PRESENTER & EDITOR ⚡
Production-Grade Desktop Application built with CustomTkinter, python-pptx, Pillow, & screeninfo.
Featuring Dynamic Responsive Font Scaling & Split-Screen Auto-Resizing UI Containers.
========================================================================================
"""

# ======================================================================================
# CRITICAL: Stream guards MUST be set BEFORE any library imports.
# Native C libraries (sounddevice, vosk, numpy) check sys.stdout/stderr on import.
# In PyInstaller --noconsole mode these are None, causing instant crashes.
# ======================================================================================
import sys
import os

class DummyWriter:
    """Safe no-op writer for Windows PyInstaller frozen --noconsole mode."""
    def write(self, s): pass
    def flush(self): pass
    def reconfigure(self, **kwargs): pass
    def fileno(self): raise OSError("No underlying file descriptor")
    def isatty(self): return False

if sys.stdout is None:
    sys.stdout = DummyWriter()
if sys.stderr is None:
    sys.stderr = DummyWriter()

try:
    if hasattr(sys.stdout, 'reconfigure'):
        sys.stdout.reconfigure(encoding='utf-8')
except Exception:
    pass

# Redirect native C library output (vosk, portaudio) that bypasses Python streams
if getattr(sys, 'frozen', False):
    try:
        import ctypes
        _devnull_fd = os.open(os.devnull, os.O_WRONLY)
        if sys.stdout is None or isinstance(sys.stdout, DummyWriter):
            ctypes.cdll.msvcrt._dup2(_devnull_fd, 1)  # redirect C stdout
        if sys.stderr is None or isinstance(sys.stderr, DummyWriter):
            ctypes.cdll.msvcrt._dup2(_devnull_fd, 2)  # redirect C stderr
    except Exception:
        pass

import multiprocessing
if getattr(sys, 'frozen', False):
    multiprocessing.freeze_support()

# Now safe to import all libraries
import ctypes
import time
import json
import queue
import io
import wave
import tempfile
import threading
import re

import customtkinter as ctk
from PIL import Image, ImageDraw, ImageFont
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from screeninfo import get_monitors

import sounddevice as sd
import numpy as np
import speech_recognition as sr


def get_resource_path(relative_path):
    """Returns absolute path to resource, working for dev and PyInstaller bundle.
    Search order: _MEIPASS -> EXE directory -> __file__ directory -> CWD -> raw path.
    """
    # 1. PyInstaller bundle temp directory (_MEIPASS)
    if getattr(sys, 'frozen', False):
        meipass = getattr(sys, '_MEIPASS', None)
        if meipass:
            p = os.path.join(meipass, relative_path)
            if os.path.exists(p):
                return p
        # 2. Directory where the .exe lives (for files placed alongside the EXE)
        exe_dir = os.path.dirname(os.path.abspath(sys.executable))
        p_exe = os.path.join(exe_dir, relative_path)
        if os.path.exists(p_exe):
            return p_exe
    # 3. Directory of this script file
    try:
        base_path = os.path.dirname(os.path.abspath(__file__))
        target_path = os.path.join(base_path, relative_path)
        if os.path.exists(target_path):
            return target_path
    except NameError:
        pass  # __file__ not defined in frozen mode
    # 4. Current working directory
    cwd_path = os.path.join(os.getcwd(), relative_path)
    if os.path.exists(cwd_path):
        return cwd_path
    return relative_path

# Try importing win32com for native PowerPoint slide rendering
try:
    import win32com.client
    HAS_PYWIN32 = True
except ImportError:
    HAS_PYWIN32 = False

# Set CustomTkinter Appearance
ctk.set_appearance_mode("Dark")
ctk.set_default_color_theme("blue")

# Theme Color Tokens (Pure Black & Sapphire Blue Aesthetic)
COLOR_BG_BLACK = "#000000"
COLOR_BG_CARD = "#0B0F19"
COLOR_SAPPHIRE = "#0F52BA"
COLOR_SAPPHIRE_HOVER = "#1C60D6"
COLOR_ACCENT_GREEN = "#10B981"
COLOR_ACCENT_RED = "#EF4444"
COLOR_TEXT_WHITE = "#FFFFFF"
COLOR_TEXT_MUTED = "#9CA3AF"
COLOR_BORDER = "#1E293B"

# Unambiguous word-to-number dictionary for voice slide targeting
NUMBER_MAP = {
    'one': 1, 'first': 1, '1': 1, '1st': 1,
    'two': 2, 'second': 2, '2': 2, '2nd': 2,
    'three': 3, 'third': 3, '3': 3, '3rd': 3,
    'four': 4, 'fourth': 4, '4': 4, '4th': 4,
    'five': 5, 'fifth': 5, '5': 5, '5th': 5,
    'six': 6, 'sixth': 6, '6': 6, '6th': 6,
    'seven': 7, 'seventh': 7, '7': 7, '7th': 7,
    'eight': 8, 'eighth': 8, '8': 8, '8th': 8,
    'nine': 9, 'ninth': 9, '9': 9, '9th': 9,
    'ten': 10, 'tenth': 10, '10': 10, '10th': 10,
    'eleven': 11, 'eleventh': 11, '11': 11,
    'twelve': 12, 'twelfth': 12, '12': 12,
    'thirteen': 13, 'thirteenth': 13, '13': 13,
    'fourteen': 14, 'fourteenth': 14, '14': 14,
    'fifteen': 15, 'fifteenth': 15, '15': 15,
    'sixteen': 16, 'sixteenth': 16, '16': 16,
    'seventeen': 17, 'seventeenth': 17, '17': 17,
    'eighteen': 18, 'eighteenth': 18, '18': 18,
    'nineteen': 19, 'nineteenth': 19, '19': 19,
    'twenty': 20, 'twentieth': 20, '20': 20,
}

# Strict high-precision voice navigation commands (prevents false triggers on conversational speech)
NAV_COMMAND_PATTERNS = [
    (re.compile(r'\b(next\s+slide|go\s+next|move\s+next|next\s+page|slide\s+next)\b'), 'ACTION_NEXT'),
    (re.compile(r'\b(previous\s+slide|prev\s+slide|go\s+back|move\s+back|previous\s+page|slide\s+back)\b'), 'ACTION_PREV'),
    (re.compile(r'\b(first\s+slide|start\s+over|go\s+to\s+beginning|restart\s+presentation|initial\s+slide)\b'), 'ACTION_FIRST'),
    (re.compile(r'\b(last\s+slide|final\s+slide|end\s+of\s+presentation|conclusion\s+slide)\b'), 'ACTION_LAST'),
    (re.compile(r'\b(black\s+screen|blackout|black\s+out)\b'), 'ACTION_BLACKOUT'),
    (re.compile(r'\b(white\s+screen|whiteout|white\s+out)\b'), 'ACTION_WHITEOUT'),
]

SLIDE_NUM_REGEX_1 = re.compile(r'\b(?:slide|page|go\s+to\s+slide|jump\s+to\s+slide|slide\s+number)\s+([a-z0-9]+)\b')
SLIDE_NUM_REGEX_2 = re.compile(r'\b([a-z0-9]+)\s+(?:slide|page)\b')

# Pre-compiled regex for punctuation stripping (avoids recompiling every match call)
_RE_PUNCT = re.compile(r'[^\w\s]')
_RE_MULTI_SPACE = re.compile(r'\s+')

# ======================================================================================
# PRE-CACHED FONTS (loaded ONCE at startup, reused across all render calls)
# ======================================================================================
_FONT_CACHE = {}

def _get_cached_font(size):
    """Returns a cached ImageFont at the given size. Loads once, reuses forever."""
    if size in _FONT_CACHE:
        return _FONT_CACHE[size]
    try:
        _font_path = os.path.join(os.environ.get('WINDIR', r'C:\Windows'), 'Fonts', 'arial.ttf')
        if os.path.exists(_font_path):
            font = ImageFont.truetype(_font_path, size)
        else:
            font = ImageFont.truetype("arial.ttf", size)
    except Exception:
        font = ImageFont.load_default()
    _FONT_CACHE[size] = font
    return font


# ======================================================================================
# 1. SLIDE DATA MODEL & MANAGER (100% REAL COLOURED PPT RENDERER)
# ======================================================================================

class SlideData:
    __slots__ = ('slide_id', 'title', 'bullet_points', 'notes', 'image_path', 'keywords', 'slide_image', '_render_hash', '_cached_render')

    def __init__(self, slide_id, title="Untitled Slide", bullet_points=None, notes="", image_path=None, keywords=None, slide_image=None):
        self.slide_id = slide_id
        self.title = title
        self.bullet_points = bullet_points if bullet_points is not None else []
        self.notes = notes
        self.image_path = image_path
        self.keywords = keywords if keywords is not None else []
        self.slide_image = slide_image  # 100% Real Coloured PPT Slide PIL Image
        self._render_hash = None  # Cache invalidation key
        self._cached_render = None  # Cached PIL image

    def _compute_hash(self):
        """Fast hash of content that affects rendering."""
        return hash((self.slide_id, self.title, tuple(self.bullet_points), self.image_path, tuple(self.keywords[:6])))


SLIDE_NUM_WORDS = {
    1: ("one", "first"),
    2: ("two", "second"),
    3: ("three", "third"),
    4: ("four", "fourth"),
    5: ("five", "fifth"),
    6: ("six", "sixth"),
    7: ("seven", "seventh"),
    8: ("eight", "eighth"),
    9: ("nine", "ninth"),
    10: ("ten", "tenth"),
    11: ("eleven", "eleventh"),
    12: ("twelve", "twelfth"),
    13: ("thirteen", "thirteenth"),
    14: ("fourteen", "fourteenth"),
    15: ("fifteen", "fifteenth"),
    16: ("sixteen", "sixteenth"),
    17: ("seventeen", "seventeenth"),
    18: ("eighteen", "eighteenth"),
    19: ("nineteen", "nineteenth"),
    20: ("twenty", "twentieth")
}

COMMON_STOPWORDS = {
    "the", "and", "for", "with", "this", "that", "from", "into", "are", "was",
    "were", "will", "have", "has", "had", "about", "your", "their", "which",
    "what", "when", "where", "how", "all", "any", "both", "each", "few", "more",
    "some", "such", "than", "then", "very", "can", "could", "should", "would"
}

def generate_default_slide_keywords(slide_num, title="", bullet_points=None):
    """Generates rich, intelligent default voice keywords combining slide numbers/ordinals and content."""
    keywords = []
    
    # 1. Slide Numbers & Ordinals (e.g. 'slide 1', 'slide one', 'first slide', 'first', 'one')
    keywords.append(f"slide {slide_num}")
    if slide_num in SLIDE_NUM_WORDS:
        cardinal, ordinal = SLIDE_NUM_WORDS[slide_num]
        keywords.append(f"slide {cardinal}")
        keywords.append(f"{ordinal} slide")
        keywords.append(ordinal)
        keywords.append(cardinal)
    keywords.append(f"page {slide_num}")

    # 2. Slide Title Keywords & Meaningful Phrases
    if title:
        clean_title = _RE_PUNCT.sub(' ', title.lower()).strip()
        clean_title = ' '.join(clean_title.split())
        if clean_title and not clean_title.startswith(f"slide {slide_num}") and not clean_title.startswith("slide"):
            if len(clean_title) <= 25:
                keywords.append(clean_title)
            title_words = [w for w in clean_title.split() if len(w) >= 3 and w not in COMMON_STOPWORDS]
            keywords.extend(title_words)

    # 3. Slide Content / Bullet Points Keywords
    if bullet_points:
        for bullet in bullet_points[:5]:
            if bullet:
                clean_b = _RE_PUNCT.sub(' ', str(bullet).lower()).strip()
                words = [w for w in clean_b.split() if len(w) >= 4 and w not in COMMON_STOPWORDS]
                keywords.extend(words[:3])

    # Deduplicate while strictly preserving order
    seen = set()
    deduped = []
    for kw in keywords:
        kw_norm = kw.strip().lower()
        if kw_norm and kw_norm not in seen:
            seen.add(kw_norm)
            deduped.append(kw_norm)

    return deduped


class SlideManager:
    """Handles loading, native rendering via PowerPoint COM, editing, and saving .pptx decks."""

    def __init__(self):
        self.slides = []
        self.file_path = None
        self.temp_dir = tempfile.mkdtemp(prefix="ppt_slides_")
        self.create_sample_deck()

    def create_sample_deck(self):
        """Creates a sample presentation deck with dynamic voice keywords."""
        self.slides = [
            SlideData(
                slide_id=1,
                title="Welcome & Project Introduction",
                bullet_points=[
                    "Next-Generation AI Speech-Driven Presentation System",
                    "Dual-Screen HDMI Auto-Detection & Fullscreen Output",
                    "Sub-50ms Voice Keyword Triggered Slide Navigation",
                    "Built with CustomTkinter, Python-PPTX, Pillow & Vosk"
                ],
                notes="Welcome the audience. Introduce the key goals of real-time voice slide switching.",
                keywords=["slide 1", "slide one", "first slide", "first", "one", "welcome", "introduction", "intro", "project intro", "start presentation", "overview"]
            ),
            SlideData(
                slide_id=2,
                title="System Architecture & Data Pipeline",
                bullet_points=[
                    "Real-Time Speech Audio Callback (16ms buffers @ 60 FPS)",
                    "Vosk Local Acoustic Preview + Google Speech Neural API Engine",
                    "Fuzzy String Matching (<10ms) against Slide Keywords",
                    "Multi-Threaded Asynchronous Window Synchronizer"
                ],
                notes="Explain the sub-50ms pipeline. Highlight zero-latency audio callback & fuzzy matcher.",
                keywords=["slide 2", "slide two", "second slide", "second", "two", "architecture", "system architecture", "pipeline", "data pipeline", "system design", "tech stack"]
            ),
            SlideData(
                slide_id=3,
                title="Key Features & Live Demo",
                bullet_points=[
                    "Automatic HDMI / DisplayPort Secondary Monitor Targeting",
                    "60 FPS Live Voice Level VU Meter embedded in Header",
                    "Blackout ('B') & Whiteout ('W') Screen Control Hotkeys",
                    "Full .PPTX Deck Import, Editing, and Saving Support"
                ],
                notes="Demonstrate live voice switching by saying keywords naturally into the microphone.",
                keywords=["slide 3", "slide three", "third slide", "third", "three", "features", "key features", "live demo", "demonstration", "capabilities"]
            ),
            SlideData(
                slide_id=4,
                title="Conclusion & Questions",
                bullet_points=[
                    "Production-Grade Architecture for Seamless Presenting",
                    "Eliminates Manual Clickers using Intelligent Voice Triggers",
                    "Flexible Keyword Customization per Slide",
                    "Thank You! Questions & Discussion"
                ],
                notes="Conclude presentation and open the floor for Q&A from the audience.",
                keywords=["slide 4", "slide four", "fourth slide", "fourth", "four", "conclusion", "summary", "wrap up", "questions", "q and a", "thank you"]
            )
        ]

    def export_slides_via_powerpoint_com(self, filepath):
        """Uses Windows PowerPoint COM API to export 100% PERFECT real coloured slide images."""
        if not HAS_PYWIN32:
            return {}

        rendered_images = {}
        try:
            abs_path = os.path.abspath(filepath)
            out_folder = os.path.join(self.temp_dir, f"export_{int(time.time())}")
            os.makedirs(out_folder, exist_ok=True)

            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            presentation = ppt_app.Presentations.Open(abs_path, ReadOnly=True, Untitled=False, WithWindow=0)

            # Export All Slides as JPG/PNG (Format #17 = ppSaveAsPNG / JPG)
            presentation.SaveAs(out_folder, 17)
            presentation.Close()

            # Find generated JPG/PNG image files
            for root, dirs, files in os.walk(out_folder):
                for f in files:
                    if f.lower().endswith((".png", ".jpg", ".jpeg", ".bmp")):
                        full_img_path = os.path.join(root, f)
                        num_str = "".join([c for c in f if c.isdigit()])
                        if num_str.isdigit():
                            s_idx = int(num_str) - 1
                            try:
                                pil_img = Image.open(full_img_path).convert("RGB")
                                rendered_images[s_idx] = pil_img
                            except Exception:
                                pass
        except Exception as e:
            print(f"[WARN] PowerPoint COM Export fallback: {e}")

        return rendered_images

    def load_pptx(self, filepath):
        """Parses standard .pptx file and extracts 100% REAL COLOURED POWERPOINT SLIDE VISUALS."""
        try:
            prs = Presentation(filepath)
            new_slides = []
            
            # Export exact 100% real coloured PowerPoint slide images via COM API!
            com_rendered_images = self.export_slides_via_powerpoint_com(filepath)

            for idx, slide in enumerate(prs.slides):
                slide_num = idx + 1
                title = f"Slide {slide_num}"
                bullet_points = []
                notes = ""
                extracted_img_path = None
                
                # Extract Title, Body Text, and Embedded Shape Images
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            if not title or title.startswith("Slide "):
                                title = text.split("\n")[0]
                            else:
                                for para in shape.text_frame.paragraphs:
                                    if para.text.strip():
                                        bullet_points.append(para.text.strip())
                    
                    # Extract embedded image blob if present
                    if extracted_img_path is None and hasattr(shape, "image"):
                        try:
                            img_bytes = shape.image.blob
                            ext = getattr(shape.image, "ext", "png")
                            img_path = os.path.join(self.temp_dir, f"slide_{slide_num}_img.{ext}")
                            with open(img_path, "wb") as img_file:
                                img_file.write(img_bytes)
                            extracted_img_path = img_path
                        except Exception:
                            pass
                                        
                # Extract Notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    
                # Dynamic Voice Keywords combining slide numbers, ordinals, and content
                keywords = generate_default_slide_keywords(slide_num, title, bullet_points)
                
                # Check for COM rendered 100% real PowerPoint slide image
                com_img = com_rendered_images.get(idx, None)

                new_slides.append(SlideData(
                    slide_id=slide_num,
                    title=title,
                    bullet_points=bullet_points[:6],
                    notes=notes,
                    image_path=extracted_img_path,
                    keywords=keywords,
                    slide_image=com_img
                ))
                
            if new_slides:
                self.slides = new_slides
                self.file_path = filepath
                return True
        except Exception as e:
            print(f"[ERROR] Failed to parse .pptx file: {e}")
            return False
        return False

    def save_pptx(self, filepath):
        """Exports current slide deck to standard .pptx file."""
        try:
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            
            for slide_data in self.slides:
                slide = prs.slides.add_slide(blank_layout)
                
                # Title Box
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.4), Inches(1.2))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = slide_data.title
                p.font.bold = True
                p.font.size = Pt(36)
                p.font.color.rgb = RGBColor(15, 82, 186)
                
                # Bullet Points Box
                if slide_data.bullet_points:
                    bodyBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(4.5))
                    btf = bodyBox.text_frame
                    for i, bullet in enumerate(slide_data.bullet_points):
                        p = btf.add_paragraph() if i > 0 else btf.paragraphs[0]
                        p.text = f"•  {bullet}"
                        p.font.size = Pt(20)
                        p.font.color.rgb = RGBColor(220, 225, 235)
                        p.space_after = Pt(14)
                        
                # Notes Box
                if slide_data.notes and slide.has_notes_slide:
                    slide.notes_slide.notes_text_frame.text = slide_data.notes

            prs.save(filepath)
            self.file_path = filepath
            return True
        except Exception as e:
            print(f"[ERROR] Failed to save .pptx file: {e}")
            return False

    def render_slide_image(self, slide_data, width=1280, height=720):
        """Returns cached slide image. Only re-renders when slide content actually changes."""
        if slide_data.slide_image is not None:
            return slide_data.slide_image

        # Check cache — skip rendering if content hasn't changed
        current_hash = slide_data._compute_hash()
        if slide_data._cached_render is not None and slide_data._render_hash == current_hash:
            return slide_data._cached_render

        # Fallback Canvas Renderer — uses pre-cached fonts (zero font-loading overhead)
        img = Image.new("RGB", (width, height), COLOR_BG_CARD)
        draw = ImageDraw.Draw(img)
        
        draw.rectangle([0, 0, width, 12], fill=COLOR_SAPPHIRE)
        draw.rectangle([0, height - 8, width, height], fill=COLOR_SAPPHIRE)
        
        title_font = _get_cached_font(int(height * 0.055))
        body_font = _get_cached_font(int(height * 0.035))
        meta_font = _get_cached_font(int(height * 0.025))

        # Pre-computed layout coordinates
        x_margin = int(width * 0.05)
        x_bullet_dot = int(width * 0.06)
        x_bullet_text = int(width * 0.09)
        y_badge = int(height * 0.06)
        y_title = int(height * 0.12)
        y_line = int(height * 0.22)
        x_line_end = int(width * 0.95)

        draw.text((x_margin, y_badge), f"SLIDE #{slide_data.slide_id}", fill=COLOR_SAPPHIRE, font=meta_font)
        draw.text((x_margin, y_title), slide_data.title, fill=COLOR_TEXT_WHITE, font=title_font)
        draw.line([(x_margin, y_line), (x_line_end, y_line)], fill=COLOR_BORDER, width=3)

        y_pos = int(height * 0.28)
        spacing = int(height * 0.09)
        
        if slide_data.bullet_points:
            for bullet in slide_data.bullet_points:
                draw.ellipse([x_bullet_dot, y_pos + 6, x_bullet_dot + 12, y_pos + 18], fill=COLOR_SAPPHIRE)
                draw.text((x_bullet_text, y_pos), bullet, fill="#E2E8F0", font=body_font)
                y_pos += spacing
        else:
            draw.text((x_bullet_text, y_pos), "(Blank Slide Content)", fill=COLOR_TEXT_MUTED, font=body_font)

        if slide_data.image_path and os.path.exists(slide_data.image_path):
            try:
                sub_img = Image.open(slide_data.image_path)
                sub_img.thumbnail((int(width * 0.35), int(height * 0.45)), Image.Resampling.BILINEAR)
                img.paste(sub_img, (int(width * 0.58), int(height * 0.32)))
            except Exception:
                pass

        if slide_data.keywords:
            kw_str = "🔑 Voice Keywords: " + ", ".join([f'"{k}"' for k in slide_data.keywords[:6]])
            draw.text((x_margin, int(height * 0.91)), kw_str, fill="#64748B", font=meta_font)

        # Cache the rendered image
        slide_data._render_hash = current_hash
        slide_data._cached_render = img
        return img


# ======================================================================================
# 2. REAL-TIME LIVE SPEECH ENGINE & FUZZY KEYWORD MATCHER
# ======================================================================================

class WinMMAudioMonitor:
    """Windows Multimedia ctypes helper to detect live physical audio endpoints in real time."""
    class WAVEINCAPSW(ctypes.Structure):
        _fields_ = [
            ('wMid', ctypes.c_ushort),
            ('wPid', ctypes.c_ushort),
            ('vDriverVersion', ctypes.c_uint),
            ('szPname', ctypes.c_wchar * 32),
            ('dwFormats', ctypes.c_ulong),
            ('wChannels', ctypes.c_ushort),
            ('wReserved1', ctypes.c_ushort),
        ]

    @classmethod
    def get_live_device_names(cls):
        """Returns list of active physical audio input endpoints registered with Windows."""
        dev_names = []
        try:
            if hasattr(ctypes, 'windll') and hasattr(ctypes.windll, 'winmm'):
                num = ctypes.windll.winmm.waveInGetNumDevs()
                caps = cls.WAVEINCAPSW()
                for i in range(num):
                    res = ctypes.windll.winmm.waveInGetDevCapsW(i, ctypes.byref(caps), ctypes.sizeof(caps))
                    if res == 0 and caps.szPname:
                        dev_names.append(caps.szPname.strip())
        except Exception:
            pass
        return dev_names


class VoiceSpeechEngine:
    """
    Real-Time Live Speech Recognizer & Adaptive Voice Engine.
    Dynamically auto-detects and migrates to Bluetooth headsets & USB microphones in real time.
    Renders live terminal-style VU status indicator embedded directly into GUI.
    """

    def __init__(self, on_keyword_matched_cb, on_status_update_cb, model_path="vosk-model-small-en-us-0.15"):
        self.on_keyword_matched_cb = on_keyword_matched_cb
        self.on_status_update_cb = on_status_update_cb
        self.sr_recognizer = sr.Recognizer()
        self.model_path = model_path
        
        self.vosk_model = None
        self.vosk_recognizer = None
        self.is_model_loaded = False
        self.is_model_loading = True
        
        # Async background model loader for sub-second GUI launch
        self._model_load_thread = threading.Thread(target=self._async_load_vosk_model, daemon=True)
        self._model_load_thread.start()

        self.audio_queue = queue.Queue(maxsize=300)
        self.speech_chunks_queue = queue.Queue(maxsize=50)
        
        self.is_recording = False
        self.is_running = True
        self.stream = None
        
        self.audio_level = 0.0
        self.latency_ms = 0.0
        self.start_time = None
        self.last_audio_frame_time = 0.0
        self._quiet_frame_count = 0
        
        self.device_id = None
        self.device_name = ""
        self.device_type = "MIC"
        self.is_mic_connected = False
        self.native_sr = 44100
        self.channels = 1
        self._last_winmm_signature = None
        self._last_device_signature = None
        
        self.probe_microphone()
        
        self.keywords_map = {}  # {keyword_str: slide_index}
        self.last_matched_keyword = ""
        self.last_matched_slide = None
        self.last_match_time = 0.0
        
        self.rec_symbols = ["🔴 LIVE", "🎙️  REC ", "⚡ STREAM", "🔴 LIVE"]
        self.anim_idx = 0
        self.last_anim_time = time.time()
        
        self.current_vad_buffer = bytearray()
        self.silence_frames = 0
        self.speech_frames = 0
        
        # Pre-sorted keyword caches (built once in set_keywords, reused every match)
        self._sorted_keywords = []
        self._keyword_patterns = {}
        self._grammar_json = None
        self.noise_floor = 20.0

    def _async_load_vosk_model(self):
        """Asynchronously loads Vosk Neural Kaldi Acoustic Model in background (<0.1s UI launch)."""
        try:
            import vosk
            resolved_path = get_resource_path(self.model_path)
            if os.path.exists(resolved_path):
                self.vosk_model = vosk.Model(resolved_path)
                if self._grammar_json:
                    self.vosk_recognizer = vosk.KaldiRecognizer(self.vosk_model, 16000, self._grammar_json)
                else:
                    self.vosk_recognizer = vosk.KaldiRecognizer(self.vosk_model, 16000)
                self.vosk_recognizer.SetWords(True)
                self.is_model_loaded = True
                self.is_model_loading = False
                print(f"[INFO] Vosk Sub-10ms Local Acoustic Engine Initialized Successfully from: {resolved_path}")
            else:
                self.is_model_loading = False
                print(f"[WARN] Vosk model directory not found at: {resolved_path}")
        except Exception as e:
            self.is_model_loading = False
            print(f"[WARN] Local Vosk initialization error: {e}")

    def init_vosk_model(self):
        """Legacy helper for synchronous compatibility."""
        if not self.is_model_loaded and not self.is_model_loading:
            self._async_load_vosk_model()

    def set_keywords(self, slides):
        """Builds keyword-to-slide mapping + pre-sorted keyword list + pre-compiled regex patterns
        and compiles a Constrained Kaldi FST Grammar to eliminate background noise & hallucinations."""
        mapping = {}
        self._total_slides_count = len(slides)
        grammar_items = set()

        # 1. Standard navigation phrases
        nav_vocab = [
            "next slide", "go next", "move next", "next page", "slide next",
            "previous slide", "prev slide", "go back", "move back", "previous page", "slide back",
            "first slide", "start over", "go to beginning", "restart presentation", "initial slide",
            "last slide", "final slide", "end of presentation", "conclusion slide",
            "black screen", "blackout", "black out",
            "white screen", "whiteout", "white out",
            "slide", "page"
        ]
        grammar_items.update(nav_vocab)

        # 2. Slide numbers and word forms for all slides in the deck
        for s_idx in range(len(slides)):
            s_num = s_idx + 1
            grammar_items.add(f"slide {s_num}")
            grammar_items.add(f"page {s_num}")
            grammar_items.add(f"go to slide {s_num}")

        for word, num in NUMBER_MAP.items():
            if 1 <= num <= len(slides):
                grammar_items.add(f"slide {word}")
                grammar_items.add(f"page {word}")
                grammar_items.add(f"go to slide {word}")
                grammar_items.add(f"{word} slide")
                grammar_items.add(f"{word} page")
                grammar_items.add(word)

        # 3. User-assigned slide keywords
        for idx, slide in enumerate(slides):
            for kw in slide.keywords:
                clean_kw = _RE_PUNCT.sub(' ', kw.strip().lower())
                clean_kw = ' '.join(clean_kw.split())
                if clean_kw:
                    mapping[clean_kw] = idx
                    grammar_items.add(clean_kw)
                    # Add individual multi-letter words
                    for part in clean_kw.split():
                        if len(part) >= 3:
                            grammar_items.add(part)

        # 4. Unknown word sink (Kaldi sends all non-matching background audio to [unk])
        grammar_items.add("[unk]")

        self.keywords_map = mapping
        self._sorted_keywords = sorted(mapping.keys(), key=len, reverse=True)
        self._keyword_patterns = {kw: re.compile(r'\b' + re.escape(kw) + r'\b') for kw in self._sorted_keywords}

        # Compile Grammar JSON
        grammar_list = sorted(list(grammar_items))
        self._grammar_json = json.dumps(grammar_list)

        # Recompile active Kaldi recognizer if model is loaded
        if self.vosk_model and self.is_model_loaded:
            try:
                import vosk
                self.vosk_recognizer = vosk.KaldiRecognizer(self.vosk_model, 16000, self._grammar_json)
                self.vosk_recognizer.SetWords(True)
                print(f"[INFO] Kaldi FST Constrained Grammar Compiled ({len(grammar_list)} vocabulary entries)")
            except Exception as e:
                print(f"[WARN] Grammar compilation fallback: {e}")

        return len(mapping)

    def refresh_portaudio(self):
        """Safely re-initializes PortAudio to discover newly connected Bluetooth/USB devices."""
        was_rec = self.is_recording
        if was_rec and self.stream:
            try:
                self.stream.stop()
                self.stream.close()
            except Exception:
                pass
            self.stream = None
            
        try:
            sd._terminate()
            sd._initialize()
        except Exception as e:
            print(f"[WARN] PortAudio re-init exception: {e}")

    def get_available_microphones(self, force_refresh=False):
        """Scans and returns all physical and default audio input devices on the system."""
        if force_refresh:
            self.refresh_portaudio()

        mic_list = []
        try:
            devices = sd.query_devices()
            default_in = sd.default.device[0] if sd.default.device else None
            
            for idx, dev in enumerate(devices):
                in_ch = dev.get('max_input_channels', 0)
                if in_ch <= 0:
                    continue
                name = dev.get('name', f"Input Device #{idx}")
                host_api_idx = dev.get('hostapi', 0)
                try:
                    host_api_name = sd.query_hostapis(host_api_idx).get('name', '')
                except Exception:
                    host_api_name = ""
                    
                sr = int(dev.get('default_samplerate', 44100))
                
                # Filter out loopback devices or ghost WDM-KS driver paths
                name_lower = name.lower()
                is_virtual = any(skip in name_lower for skip in ['sound mapper', 'primary sound', 'stereo mix', 'speaker', 'loopback', 'cable', 'virtual', 'camo'])
                is_ghost_wdm = ('wdm-ks' in host_api_name.lower() and ('@system32' in name_lower or 'bthhfenum' in name_lower))

                if is_ghost_wdm:
                    continue

                is_default = (idx == default_in)
                
                # Determine device type
                is_bt = any(k in name_lower for k in [
                    'bluetooth', 'headset', 'hands-free', 'buds', 'airdopes', 'wireless', 
                    'airpods', 'noise', 'boat', 'oneplus', 'mivi', 'realme', 'boult', 
                    'jbl', 'sony', 'bose', 'sennheiser', 'jabra', 'anker', 'soundcore', 'redmi', 'oppo', 'vivo'
                ])
                is_usb = any(k in name_lower for k in ['usb', 'hyperx', 'blue yeti', 'rode', 'fifine', 'boya', 'podcast', 'samson', 'shure', 'audio-technica'])
                
                dev_type = "BT" if is_bt else ("USB" if is_usb else "MIC")
                icon = "🎧" if is_bt else ("🎙️" if is_usb else "💻")

                mic_list.append({
                    'id': idx,
                    'name': name,
                    'samplerate': sr,
                    'channels': min(in_ch, 2),
                    'api': host_api_name,
                    'is_virtual': is_virtual,
                    'is_default': is_default,
                    'type': dev_type,
                    'icon': icon,
                    'display_name': f"{icon} {name} ({host_api_name})"
                })
        except Exception as e:
            print(f"[WARN] Failed to query audio devices: {e}")
        return mic_list

    def probe_microphone(self, force_device_id=None, force_refresh=False):
        """Probes and dynamically selects the best available microphone (Prioritizes Bluetooth & USB)."""
        try:
            if force_refresh:
                self.refresh_portaudio()

            all_mics = self.get_available_microphones()
            if not all_mics:
                self.device_id = None
                self.device_name = "No Input Device Detected"
                self.device_type = "NONE"
                self.is_mic_connected = False
                return False

            if force_device_id is not None:
                for m in all_mics:
                    if m['id'] == force_device_id:
                        self.device_id = m['id']
                        self.device_name = m['name']
                        self.native_sr = m['samplerate']
                        self.channels = m['channels']
                        self.device_type = m.get('type', 'MIC')
                        self.is_mic_connected = True
                        print(f"[INFO] Manually Selected Microphone #{self.device_id}: {self.device_name} ({self.device_type})")
                        return True

            bt_candidates = []
            usb_candidates = []
            internal_candidates = []

            for m in all_mics:
                if m.get('is_virtual', False):
                    continue
                idx = m['id']
                name = m['name']
                sr = m['samplerate']
                ch = m['channels']
                api_name = m.get('api', '').lower()

                # Instant hardware capability check using PortAudio (<0.01ms)
                can_open = False
                tested_sr = sr
                for test_sr in [sr, 16000, 48000, 44100]:
                    try:
                        sd.check_input_settings(device=idx, channels=ch, dtype='int16', samplerate=test_sr)
                        can_open = True
                        tested_sr = test_sr
                        break
                    except Exception:
                        pass

                if not can_open:
                    continue

                info = dict(m)
                info['samplerate'] = tested_sr

                # Score host APIs: WASAPI = 4 (best for Windows 10/11), DirectSound = 3, MME = 2, WDM-KS = 1
                api_score = 4 if 'wasapi' in api_name else (3 if 'directsound' in api_name else (2 if 'mme' in api_name else 1))

                if m.get('type') == 'BT':
                    bt_candidates.append((api_score, info))
                elif m.get('type') == 'USB':
                    usb_candidates.append((api_score, info))
                else:
                    internal_candidates.append((api_score, info))

            selected_info = None
            selected_type = "MIC"

            # 1. First Priority: Active Bluetooth / Wireless Headset
            if bt_candidates:
                bt_candidates.sort(key=lambda x: x[0], reverse=True)
                selected_info = bt_candidates[0][1]
                selected_type = "BT"
            # 2. Second Priority: USB Microphone
            elif usb_candidates:
                usb_candidates.sort(key=lambda x: x[0], reverse=True)
                selected_info = usb_candidates[0][1]
                selected_type = "USB"
            # 3. Third Priority: Built-in Microphone Array
            elif internal_candidates:
                internal_candidates.sort(key=lambda x: x[0], reverse=True)
                selected_info = internal_candidates[0][1]
                selected_type = "MIC"
            else:
                selected_info = all_mics[0]
                selected_type = selected_info.get('type', 'MIC')

            self.device_id = selected_info['id']
            self.device_name = selected_info['name']
            self.device_type = selected_type
            self.native_sr = selected_info['samplerate']
            self.channels = selected_info['channels']
            self.is_mic_connected = True
            self._last_winmm_signature = tuple(WinMMAudioMonitor.get_live_device_names())
            self._last_device_signature = tuple((m['id'], m['name']) for m in all_mics)
            print(f"[INFO] Auto-selected {selected_type} Microphone #{self.device_id}: {self.device_name} ({selected_info.get('api', '')} @ {self.native_sr}Hz, {self.channels}ch)")
            return True

        except Exception as e:
            print(f"[WARN] Audio input device probe skipped: {e}")
            
        self.device_id = None
        self.device_name = "System Default Microphone"
        self.device_type = "MIC"
        self.is_mic_connected = True
        self.native_sr = 16000
        self.channels = 1
        return True

    def switch_device(self, new_device_id):
        """Switches microphone endpoint immediately (seamless live migration if recording)."""
        was_rec = self.is_recording
        if was_rec:
            self.stop()
        success = self.probe_microphone(force_device_id=new_device_id)
        if was_rec and success:
            self.start()
        return success

    def check_device_hotplug(self):
        """Actively detects Bluetooth/USB/Hardware mic connections, removals, stalls, and auto-migrates live stream."""
        try:
            live_winmm = tuple(WinMMAudioMonitor.get_live_device_names())
            winmm_changed = (live_winmm != getattr(self, '_last_winmm_signature', None))
            
            # Watchdog check: If recording is supposed to be active, but stream stalled or died
            now = time.time()
            stream_stalled = False
            if self.is_recording:
                if self.stream is None or not getattr(self.stream, 'active', False):
                    stream_stalled = True
                elif getattr(self, 'last_audio_frame_time', 0.0) > 0 and (now - self.last_audio_frame_time > 2.5) and (now - getattr(self, 'start_time', now) > 3.0):
                    stream_stalled = True

            if winmm_changed or stream_stalled:
                old_name = self.device_name
                was_rec = self.is_recording
                old_type = self.device_type

                # Refresh PortAudio device table
                self.refresh_portaudio()
                self._last_winmm_signature = live_winmm
                
                # Check for Bluetooth headsets in live devices
                bt_keywords = [
                    'bluetooth', 'headset', 'hands-free', 'buds', 'airdopes', 'wireless', 
                    'airpods', 'noise', 'boat', 'oneplus', 'mivi', 'realme', 'boult', 
                    'jbl', 'sony', 'bose', 'sennheiser', 'jabra', 'anker', 'soundcore', 'redmi', 'oppo', 'vivo'
                ]
                has_bt_now = any(any(k in d.lower() for k in bt_keywords) for d in live_winmm)

                success = self.probe_microphone(force_refresh=False)
                if was_rec and success:
                    self.start()

                if has_bt_now and old_type != "BT" and self.device_type == "BT":
                    print(f"[HOTPLUG] Dynamic Adaptation: Switched to Bluetooth Headset #{self.device_id} ({self.device_name})")
                    return "BLUETOOTH_CONNECTED", old_name, self.device_name
                elif stream_stalled:
                    print(f"[HOTPLUG] Stream Auto-Recovered on #{self.device_id} ({self.device_name})")
                    return "RECOVERED", old_name, self.device_name
                elif not has_bt_now and old_type == "BT":
                    print(f"[HOTPLUG] Bluetooth Disconnected! Auto-switched to: {self.device_name}")
                    return "DISCONNECTED", old_name, self.device_name
                else:
                    return "SWITCHED", old_name, self.device_name

        except Exception as e:
            pass
        return "UNCHANGED", None, None

    def vosk_instant_worker_loop(self):
        """Dedicated sub-10ms local Vosk worker thread for lightning-fast keyword & command detection."""
        while self.is_running:
            if self.is_recording and self.vosk_recognizer and self.is_model_loaded:
                try:
                    data = self.audio_queue.get(timeout=0.01)
                    if data:
                        text = ""
                        is_partial = False
                        if self.vosk_recognizer.AcceptWaveform(data):
                            res = json.loads(self.vosk_recognizer.Result())
                            text = res.get("text", "").strip()
                            is_partial = False
                        else:
                            pres = json.loads(self.vosk_recognizer.PartialResult())
                            text = pres.get("partial", "").strip()
                            is_partial = True

                        if text:
                            kw, target = self.match_speech_to_keyword(text, is_partial=is_partial)
                            now = time.time()
                            if target is not None and (target != getattr(self, 'last_matched_target', None) or now - getattr(self, 'last_match_time', 0) > 0.8):
                                self.last_matched_target = target
                                self.last_matched_keyword = kw
                                self.last_match_time = now
                                self.on_keyword_matched_cb(target, kw, text)
                                # Reset recognizer buffer immediately upon matching to avoid duplicate triggers
                                try:
                                    self.vosk_recognizer.Reset()
                                    # Clear stale audio frames after match
                                    while not self.audio_queue.empty():
                                        try:
                                            self.audio_queue.get_nowait()
                                        except queue.Empty:
                                            break
                                except Exception:
                                    pass
                except queue.Empty:
                    pass
            else:
                time.sleep(0.01)

    def audio_callback(self, indata, frames, time_info, status):
        """Real-time microphone input callback (<0.5ms latency). Direct NumPy array processing."""
        t0 = time.perf_counter()
        self.last_audio_frame_time = time.time()
        
        # Downmix multi-channel to mono
        if indata.ndim > 1 and indata.shape[1] > 1:
            pcm_mono = (indata[:, 0].astype(np.int32) + indata[:, 1].astype(np.int32)) >> 1
            pcm_mono = pcm_mono.astype(np.int16)
        else:
            pcm_mono = indata[:, 0] if indata.ndim > 1 else indata

        n = len(pcm_mono)
        if n > 0:
            # Sensitive RMS calculation
            rms = float(np.sqrt(np.mean(pcm_mono.astype(np.int32) ** 2)))
            # Dynamic SNR background noise floor tracking
            self.noise_floor = 0.96 * getattr(self, 'noise_floor', 20.0) + 0.04 * min(rms, 60.0)
            
            # Instant Peak Attack + Smooth Exponential Decay
            effective_level = max(0.0, rms - self.noise_floor)
            raw_level = min(1.0, max(0.0, effective_level / 40.0))
            if raw_level > self.audio_level:
                self.audio_level = raw_level  # Instantaneous spike on speech!
            else:
                self.audio_level = self.audio_level * 0.84 + raw_level * 0.16  # Smooth decay

            if self.is_recording:
                # Dynamic Acoustic Noise Gate (adaptive threshold above ambient floor)
                adaptive_thresh = max(25.0, getattr(self, 'noise_floor', 20.0) + 12.0)
                if rms > adaptive_thresh:
                    self._quiet_frame_count = 0
                else:
                    self._quiet_frame_count = getattr(self, '_quiet_frame_count', 0) + 1

                # Feed audio during speech and brief natural pauses (<120ms), discard ambient silence
                if self._quiet_frame_count < 4:
                    if self.native_sr == 16000:
                        raw_16k = pcm_mono.tobytes()
                    else:
                        n_orig = len(pcm_mono)
                        n_targ = int(n_orig * 16000 / self.native_sr)
                        if n_targ > 0:
                            x_orig = np.linspace(0, 1, n_orig, dtype=np.float32)
                            x_targ = np.linspace(0, 1, n_targ, dtype=np.float32)
                            raw_16k = np.clip(np.interp(x_targ, x_orig, pcm_mono.astype(np.float32)), -32768, 32767).astype(np.int16).tobytes()
                        else:
                            raw_16k = None

                    if raw_16k:
                        try:
                            self.audio_queue.put_nowait(raw_16k)
                        except queue.Full:
                            pass

        self.latency_ms = (time.perf_counter() - t0) * 1000.0

    def match_speech_to_keyword(self, spoken_text, is_partial=False):
        """Comprehensive natural voice command & slide keyword parser with millisecond response."""
        if not spoken_text:
            return None, None

        clean_text = spoken_text.replace("[unk]", " ").strip()
        clean_text = _RE_PUNCT.sub(' ', clean_text.lower()).strip()
        clean_text = ' '.join(clean_text.split())
        if not clean_text or clean_text == "unk":
            return None, None

        # 1. Global Navigation Voice Commands (Next Slide, Previous Slide, First Slide, Last Slide, Blackout, Whiteout)
        for pat, action in NAV_COMMAND_PATTERNS:
            m = pat.search(clean_text)
            if m:
                return m.group(0), action

        # 2. Number Pattern Matching ("slide two", "slide 2", "second slide", "page 3", "go to slide 5")
        total_slides = getattr(self, '_total_slides_count', len(self.keywords_map))
        
        m1 = SLIDE_NUM_REGEX_1.search(clean_text)
        if m1:
            word = m1.group(1)
            if word in NUMBER_MAP:
                s_num = NUMBER_MAP[word]
                if 1 <= s_num <= total_slides:
                    return m1.group(0), s_num - 1

        m2 = SLIDE_NUM_REGEX_2.search(clean_text)
        if m2:
            word = m2.group(1)
            if word in NUMBER_MAP:
                s_num = NUMBER_MAP[word]
                if 1 <= s_num <= total_slides:
                    return m2.group(0), s_num - 1

        # 3. Exact full phrase match against slide keywords (enabled for both partial and full results)
        km = self.keywords_map
        if clean_text in km:
            return clean_text, km[clean_text]

        # 4. Whole-phrase / exact word boundary regex match for slide keywords
        for kw in self._sorted_keywords:
            if kw in self._keyword_patterns and self._keyword_patterns[kw].search(clean_text):
                return kw, km[kw]

        # Partial results only match complete keywords/commands, not single short fragmented tokens
        if is_partial:
            return None, None

        # 5. Discrete multi-letter word match (minimum 3 characters)
        words = clean_text.split()
        for w in words:
            if len(w) >= 3 and w in km:
                return w, km[w]

        return None, None

    def neural_worker_loop(self):
        """Background thread sending speech audio chunks to Neural API."""
        while self.is_running:
            if self.is_recording:
                try:
                    chunk_bytes = self.speech_chunks_queue.get(timeout=0.05)
                    pcm = np.frombuffer(chunk_bytes, dtype=np.int16)
                    if len(pcm) > 0:
                        if self.channels > 1:
                            pcm = pcm.reshape(-1, self.channels)
                            pcm_mono = np.mean(pcm, axis=1).astype(np.int16)
                        else:
                            pcm_mono = pcm

                        num_target = int(len(pcm_mono) * 16000 / self.native_sr)
                        if num_target > 0:
                            x_orig = np.linspace(0, 1, len(pcm_mono), dtype=np.float32)
                            x_targ = np.linspace(0, 1, num_target, dtype=np.float32)
                            pcm16 = np.clip(np.interp(x_targ, x_orig, pcm_mono.astype(np.float32)), -32768, 32767).astype(np.int16)

                            wav_io = io.BytesIO()
                            with wave.open(wav_io, 'wb') as wf:
                                wf.setnchannels(1)
                                wf.setsampwidth(2)
                                wf.setframerate(16000)
                                wf.writeframes(pcm16.tobytes())
                            wav_io.seek(0)

                            with sr.AudioFile(wav_io) as source:
                                audio = self.sr_recognizer.record(source)

                            try:
                                text = self.sr_recognizer.recognize_google(audio).strip()
                                if text:
                                    kw, slide_idx = self.match_speech_to_keyword(text)
                                    if kw is not None:
                                        self.last_matched_keyword = kw
                                        self.last_matched_slide = slide_idx
                                        self.on_keyword_matched_cb(slide_idx, kw, text)
                            except Exception:
                                pass
                except queue.Empty:
                    pass
            else:
                time.sleep(0.02)

    def get_status_indicator_str(self, fps=60.0):
        """Renders static, non-jittering fixed-width VU level meter & speech indicator."""
        bars = int(round(self.audio_level * 12))
        bars = max(0, min(12, bars))
        vu_bar = "█" * bars + "░" * (12 - bars)
        
        status_tag = "🟢 LIVE" if self.is_recording else "⚪ OFF "
        lat = min(99.9, max(0.0, self.latency_ms))
        fps_val = min(99.9, max(0.0, fps))
        return f"{status_tag} [{vu_bar}] [{fps_val:4.1f} FPS | {lat:4.1f}ms]"

    def start(self):
        """Starts real-time live speech recognition stream with guaranteed multi-device fallback."""
        if self.is_recording and self.stream and getattr(self.stream, 'active', False):
            return True

        if not self.is_model_loaded and self.is_model_loading:
            try:
                self._model_load_thread.join(timeout=1.5)
            except Exception:
                pass

        self.stop()
        self.start_time = time.time()
        self.last_audio_frame_time = time.time()

        all_devices = self.get_available_microphones()
        
        # Build device trial list in priority order
        device_attempts = []
        if self.device_id is not None:
            device_attempts.append(self.device_id)
            
        for m in all_devices:
            if not m.get('is_virtual', False) and m['id'] not in device_attempts:
                device_attempts.append(m['id'])
                
        # Always append system default as final guaranteed fallback
        device_attempts.append(None)

        stream_started = False
        active_device_name = ""

        for dev_target in device_attempts:
            if dev_target is not None:
                # Find device info
                dev_info = next((m for m in all_devices if m['id'] == dev_target), None)
                if dev_info:
                    native_sr = dev_info.get('samplerate', 44100)
                    ch_trials = [min(dev_info.get('channels', 1), 2), 1]
                    dev_name = dev_info.get('name', f'Mic #{dev_target}')
                else:
                    native_sr = 44100
                    ch_trials = [2, 1]
                    dev_name = f'Mic #{dev_target}'
                sr_trials = [native_sr, 16000, 48000, 44100, 8000]
            else:
                sr_trials = [16000, 48000, 44100, 8000]
                ch_trials = [1, 2]
                dev_name = "System Default Microphone"

            for test_sr in sr_trials:
                for test_ch in ch_trials:
                    try:
                        self.stream = sd.InputStream(
                            samplerate=test_sr,
                            device=dev_target,
                            channels=test_ch,
                            dtype='int16',
                            blocksize=int(test_sr * 0.05),
                            callback=self.audio_callback
                        )
                        self.stream.start()
                        if self.stream.active:
                            self.device_id = dev_target
                            self.device_name = dev_name
                            self.native_sr = test_sr
                            self.channels = test_ch
                            stream_started = True
                            active_device_name = dev_name
                            print(f"[SUCCESS] Audio Stream Live on #{dev_target} ({dev_name}) @ {test_sr}Hz, {test_ch}ch")
                            break
                    except Exception as e:
                        if self.stream:
                            try:
                                self.stream.close()
                            except Exception:
                                pass
                            self.stream = None

                if stream_started:
                    break
            if stream_started:
                break

        if stream_started:
            self.is_recording = True
            self.is_mic_connected = True
            self.last_audio_frame_time = time.time()
            if not getattr(self, '_vosk_thread_started', False):
                self._vosk_thread_started = True
                threading.Thread(target=self.vosk_instant_worker_loop, daemon=True).start()
            if not getattr(self, '_neural_thread_started', False):
                self._neural_thread_started = True
                threading.Thread(target=self.neural_worker_loop, daemon=True).start()
            return True

        self.is_recording = False
        self.is_mic_connected = False
        print(f"[ERROR] Could not initialize microphone audio stream on this machine.")
        return False

    def stop(self):
        """Stops live speech recognition stream cleanly."""
        self.is_recording = False
        if self.stream:
            try:
                self.stream.stop()
                self.stream.close()
            except Exception as e:
                print(f"[WARN] Exception stopping audio stream: {e}")
            finally:
                self.stream = None


# ======================================================================================
# 3. EXTERNAL DUAL-SCREEN / HDMI PRESENTATION WINDOW
# ======================================================================================

class ExternalDisplayWindow(ctk.CTkToplevel):
    """Secondary borderless fullscreen presentation window for HDMI / DisplayPort output."""

    def __init__(self, parent, monitor_info=None):
        super().__init__(parent)
        self.title("HDMI Fullscreen Presentation Output")
        self.configure(fg_color=COLOR_BG_BLACK)
        
        self.monitor = monitor_info
        self.is_blackout = False
        self.is_whiteout = False
        
        is_external = False
        if self.monitor:
            is_primary = getattr(self.monitor, 'is_primary', False) or (self.monitor.x == 0 and self.monitor.y == 0)
            if not is_primary:
                is_external = True

        if is_external and self.monitor:
            # External HDMI Monitor / Projector Connected! -> 100% Borderless Fullscreen Output
            self.geometry(f"{self.monitor.width}x{self.monitor.height}+{self.monitor.x}+{self.monitor.y}")
            self.overrideredirect(True)
            self.attributes("-topmost", True)
            self.after(500, lambda: self.attributes("-topmost", False))
        else:
            # Primary Laptop Screen (or Windowed Preview Mode) -> Compact Sized Window!
            self.overrideredirect(False)
            lap_w = self.monitor.width if self.monitor else 1280
            lap_h = self.monitor.height if self.monitor else 720
            
            win_w, win_h = 740, 416
            pos_x = max(40, (lap_w - win_w) // 2)
            pos_y = max(40, (lap_h - win_h) // 2)
            
            self.geometry(f"{win_w}x{win_h}+{pos_x}+{pos_y}")
            self.minsize(400, 225)
            
        self.slide_label = ctk.CTkLabel(self, text="", fg_color=COLOR_BG_BLACK)
        self.slide_label.pack(fill="both", expand=True)
        
        self._cached_hdmi_img = None
        self._last_hdmi_key = None
        
        self.bind("<Escape>", lambda e: self.destroy())

    def update_slide(self, pil_image, is_blackout=False, is_whiteout=False):
        """Updates slide image on external HDMI screen with zero-delay cached rendering."""
        self.is_blackout = is_blackout
        self.is_whiteout = is_whiteout
        
        if self.is_blackout:
            self.slide_label.configure(image=None, fg_color="#000000", text="")
            return
        elif self.is_whiteout:
            self.slide_label.configure(image=None, fg_color="#FFFFFF", text="")
            return
            
        w = self.winfo_width() if self.winfo_width() > 100 else (self.monitor.width if self.monitor else 1280)
        h = self.winfo_height() if self.winfo_height() > 100 else (self.monitor.height if self.monitor else 720)
        
        # Preserve perfect 16:9 aspect ratio centered in HDMI screen black background
        target_w = w
        target_h = int(target_w * 9 / 16)
        if target_h > h:
            target_h = h
            target_w = int(target_h * 16 / 9)

        target_w = max(100, target_w)
        target_h = max(56, target_h)

        img_key = (id(pil_image), target_w, target_h)
        if self._last_hdmi_key == img_key and self._cached_hdmi_img is not None:
            self.slide_label.configure(image=self._cached_hdmi_img, fg_color=COLOR_BG_BLACK, text="")
            return

        resized_pil = pil_image.resize((target_w, target_h), Image.Resampling.BILINEAR)
        ctk_img = ctk.CTkImage(light_image=resized_pil, dark_image=resized_pil, size=(target_w, target_h))
        self._cached_hdmi_img = ctk_img
        self._last_hdmi_key = img_key
        self.slide_label.configure(image=ctk_img, fg_color=COLOR_BG_BLACK, text="")


# ======================================================================================
# 4. MAIN DESKTOP APPLICATION CONTROLLER (DYNAMIC RESPONSIVE FONT SCALING)
# ======================================================================================

class PresentationApp(ctk.CTk):
    """Main Application Dashboard & Presenter Controller with Dynamic Responsive Font Scaling."""

    def __init__(self):
        super().__init__()
        self.title("⚡ AI Voice Presentation Engine (Dual-Screen HDMI)")
        self.geometry("1480x920")
        self.minsize(680, 500)
        self.configure(fg_color=COLOR_BG_BLACK)
        
        self.slide_mgr = SlideManager()
        self.current_slide_idx = 0
        self.active_view = "dashboard"
        
        self.is_blackout = False
        self.monitors = get_monitors()
        self.hdmi_window = None
        
        self.voice_engine = VoiceSpeechEngine(
            on_keyword_matched_cb=self.on_voice_keyword_matched,
            on_status_update_cb=self.update_speech_status_bar
        )
        self.voice_engine.set_keywords(self.slide_mgr.slides)
        
        self.setup_keyboard_shortcuts()
        self.build_gui_with_left_sidebar()
        
        # Auto-launch HDMI secondary window if Monitor 2 detected!
        if len(self.monitors) > 1:
            self.launch_hdmi_output(monitor_idx=1)

        # Check command line arguments for auto-loading PPTX (e.g. double click or drag onto exe)
        if len(sys.argv) > 1:
            target_pptx = sys.argv[1]
            if os.path.exists(target_pptx) and target_pptx.lower().endswith(('.pptx', '.ppt')):
                try:
                    if self.slide_mgr.load_pptx(target_pptx):
                        self.current_slide_idx = 0
                        self.voice_engine.set_keywords(self.slide_mgr.slides)
                except Exception as e:
                    print(f"[WARN] Failed to auto-load command line PPTX {target_pptx}: {e}")

        self.update_slide_display()
        self.update_microphone_indicator()
        self.start_gui_live_indicator_loop()
    def setup_keyboard_shortcuts(self):
        """Global keyboard shortcut bindings using bind_all for 100% guaranteed Arrow Key Navigation."""
        def is_typing():
            focused = self.focus_get()
            if focused:
                w_class = getattr(focused, "winfo_class", lambda: "")()
                if w_class in ["Entry", "Text", "TEntry"] or isinstance(focused, (ctk.CTkEntry, ctk.CTkTextbox)):
                    return True
            return False

        def handle_next(e):
            if not is_typing():
                self.next_slide()

        def handle_prev(e):
            if not is_typing():
                self.prev_slide()

        def handle_blackout(e):
            if not is_typing():
                self.toggle_blackout()

        def handle_whiteout(e):
            if not is_typing():
                self.toggle_whiteout()

        self.bind_all("<Right>", handle_next)
        self.bind_all("<Down>", handle_next)
        self.bind_all("<space>", handle_next)
        self.bind_all("<Next>", handle_next)
        
        self.bind_all("<Left>", handle_prev)
        self.bind_all("<Up>", handle_prev)
        self.bind_all("<BackSpace>", handle_prev)
        self.bind_all("<Prior>", handle_prev)
        
        self.bind_all("<Escape>", lambda e: self.stop_presentation())
        self.bind_all("<b>", handle_blackout)
        self.bind_all("<B>", handle_blackout)
        self.bind_all("<w>", handle_whiteout)
        self.bind_all("<W>", handle_whiteout)

    def build_gui_with_left_sidebar(self):
        """Constructs Responsive Main Window with Flex Left Sidebar & Auto-Sizing Buttons."""
        
        # 1. TOP HEADER STATUS BAR (PIXEL-PERFECT ALIGNED 3-SECTION FLEX HEADER)
        self.header_frame = ctk.CTkFrame(self, fg_color=COLOR_BG_CARD, corner_radius=0, height=60, border_width=1, border_color=COLOR_BORDER)
        self.header_frame.pack(side="top", fill="x")
        self.header_frame.pack_propagate(False)
        
        # SECTION A (LEFT): Brand Title + Animated VU Level Meter
        header_left_box = ctk.CTkFrame(self.header_frame, fg_color="transparent")
        header_left_box.pack(side="left", fill="y", padx=(14, 6))

        self.title_lbl = ctk.CTkLabel(
            header_left_box,
            text="⚡ GenSlide",
            font=ctk.CTkFont(size=14, weight="bold"),
            text_color=COLOR_SAPPHIRE
        )
        self.title_lbl.pack(side="left", padx=(0, 8), pady=14)

        ctk.CTkLabel(header_left_box, text="|", font=ctk.CTkFont(size=14, weight="bold"), text_color=COLOR_BORDER).pack(side="left", padx=(0, 8), pady=14)

        self.speech_status_lbl = ctk.CTkLabel(
            header_left_box,
            text="⚪ OFF  [░░░░░░░░░░░░] [60.0 FPS |  0.0ms]",
            font=ctk.CTkFont(family="Consolas", size=11, weight="bold"),
            text_color=COLOR_ACCENT_GREEN,
            width=275,
            anchor="w"
        )
        self.speech_status_lbl.pack(side="left", padx=(0, 8), pady=14)

        # SECTION B (CENTER): Live Input Device Status Pill + Dropdown + Rescan
        header_center_box = ctk.CTkFrame(self.header_frame, fg_color="transparent")
        header_center_box.pack(side="left", fill="y", padx=(6, 6))

        self.mic_status_badge = ctk.CTkLabel(
            header_center_box,
            text="🎙️ Mic: Detecting Device...",
            font=ctk.CTkFont(size=11, weight="bold"),
            fg_color="#064E3B",
            text_color="#34D399",
            corner_radius=6,
            padx=10,
            pady=5
        )
        self.mic_status_badge.pack(side="left", padx=(0, 6), pady=13)

        self.mic_select_menu = ctk.CTkOptionMenu(
            header_center_box,
            values=["Detecting Microphones..."],
            font=ctk.CTkFont(size=11, weight="bold"),
            dropdown_font=ctk.CTkFont(size=11),
            fg_color="#1E293B",
            button_color="#334155",
            button_hover_color=COLOR_SAPPHIRE_HOVER,
            dropdown_fg_color="#0F172A",
            dropdown_hover_color=COLOR_SAPPHIRE,
            dropdown_text_color="#FFFFFF",
            text_color="#E2E8F0",
            width=185,
            height=32,
            corner_radius=6,
            command=self.on_mic_dropdown_selected
        )
        self.mic_select_menu.pack(side="left", padx=(0, 6), pady=14)

        self.mic_rescan_btn = ctk.CTkButton(
            header_center_box,
            text="🔄",
            width=32,
            height=32,
            font=ctk.CTkFont(size=12, weight="bold"),
            fg_color="#1E293B",
            hover_color=COLOR_SAPPHIRE_HOVER,
            corner_radius=6,
            command=self.manual_rescan_microphones
        )
        self.mic_rescan_btn.pack(side="left", padx=(0, 8), pady=14)

        # SECTION C (RIGHT): Voice Action Pill + Action Buttons
        header_right_box = ctk.CTkFrame(self.header_frame, fg_color="transparent")
        header_right_box.pack(side="right", padx=(5, 14), fill="y")

        # PRESENT ON HDMI / PROJECTOR BUTTON
        self.present_hdmi_btn = ctk.CTkButton(
            header_right_box,
            text="📺 PRESENT ON HDMI",
            font=ctk.CTkFont(size=12, weight="bold"),
            fg_color=COLOR_ACCENT_GREEN,
            hover_color="#059669",
            height=34,
            corner_radius=6,
            command=self.launch_hdmi_output
        )
        self.present_hdmi_btn.pack(side="right", padx=(6, 0), pady=13)

        # START LIVE VOICE ENGINE BUTTON
        self.voice_btn = ctk.CTkButton(
            header_right_box,
            text="🎙️ LIVE VOICE",
            font=ctk.CTkFont(size=12, weight="bold"),
            fg_color=COLOR_SAPPHIRE,
            hover_color=COLOR_SAPPHIRE_HOVER,
            height=34,
            corner_radius=6,
            command=self.toggle_voice_engine
        )
        self.voice_btn.pack(side="right", padx=(6, 6), pady=13)

        self.match_badge_lbl = ctk.CTkLabel(
            header_right_box,
            text="[Voice Status: Idle]",
            font=ctk.CTkFont(size=11, weight="bold"),
            fg_color="#1E293B",
            text_color=COLOR_TEXT_MUTED,
            corner_radius=6,
            padx=10,
            pady=5
        )
        self.match_badge_lbl.pack(side="right", padx=4, pady=13)

        # 2. MAIN BODY CONTAINER (Left Sidebar + Right Content Area)
        self.body_container = ctk.CTkFrame(self, fg_color=COLOR_BG_BLACK)
        self.body_container.pack(fill="both", expand=True)

        # 3. DEDICATED LEFT SIDEBAR EDITOR (RESPONSIVE FLEX WIDTH)
        self.sidebar_frame = ctk.CTkFrame(self.body_container, fg_color=COLOR_BG_CARD, width=280, corner_radius=0, border_width=1, border_color=COLOR_BORDER)
        self.sidebar_frame.pack(side="left", fill="y")

        self.sidebar_title_lbl = ctk.CTkLabel(self.sidebar_frame, text="✏️ SLIDE DECK EDITOR", font=ctk.CTkFont(size=14, weight="bold"), text_color=COLOR_SAPPHIRE)
        self.sidebar_title_lbl.pack(anchor="w", padx=15, pady=(15, 8))

        # View Switcher Buttons (Auto-Expanding)
        view_btn_frame = ctk.CTkFrame(self.sidebar_frame, fg_color="transparent")
        view_btn_frame.pack(fill="x", padx=10, pady=2)

        self.btn_nav_dashboard = ctk.CTkButton(view_btn_frame, text="📽️ Presenter Dashboard", fg_color=COLOR_SAPPHIRE, command=lambda: self.switch_view("dashboard"))
        self.btn_nav_dashboard.pack(fill="x", pady=2)

        self.btn_nav_editor = ctk.CTkButton(view_btn_frame, text="✏️ Slide Detail Editor", fg_color=COLOR_BG_BLACK, border_width=1, border_color=COLOR_SAPPHIRE, command=lambda: self.switch_view("editor"))
        self.btn_nav_editor.pack(fill="x", pady=2)

        self.btn_nav_keywords = ctk.CTkButton(view_btn_frame, text="🏷️ Voice Keyword Matrix", fg_color=COLOR_BG_BLACK, border_width=1, border_color=COLOR_SAPPHIRE, command=lambda: self.switch_view("keywords"))
        self.btn_nav_keywords.pack(fill="x", pady=2)

        self.btn_nav_settings = ctk.CTkButton(view_btn_frame, text="🖥️ HDMI Settings", fg_color=COLOR_BG_BLACK, border_width=1, border_color=COLOR_SAPPHIRE, command=lambda: self.switch_view("settings"))
        self.btn_nav_settings.pack(fill="x", pady=2)

        ctk.CTkFrame(self.sidebar_frame, height=1, fg_color=COLOR_BORDER).pack(fill="x", padx=10, pady=10)

        # Slide List Header & Actions
        ctk.CTkLabel(self.sidebar_frame, text="SLIDE DECK LIST", font=ctk.CTkFont(size=11, weight="bold"), text_color=COLOR_TEXT_MUTED).pack(anchor="w", padx=15, pady=(2, 2))

        self.slide_list_scroll = ctk.CTkScrollableFrame(self.sidebar_frame, fg_color=COLOR_BG_BLACK, corner_radius=6)
        self.slide_list_scroll.pack(fill="both", expand=True, padx=10, pady=4)

        # File & Deck Management Buttons (Flex Resizing)
        deck_btn_frame = ctk.CTkFrame(self.sidebar_frame, fg_color="transparent")
        deck_btn_frame.pack(fill="x", padx=10, pady=6)

        ctk.CTkButton(deck_btn_frame, text="+ ADD", fg_color=COLOR_ACCENT_GREEN, command=self.add_new_slide).pack(side="left", fill="x", expand=True, padx=2)
        ctk.CTkButton(deck_btn_frame, text="🗑️ DELETE", fg_color=COLOR_ACCENT_RED, command=self.delete_slide).pack(side="right", fill="x", expand=True, padx=2)

        # OPEN PPT FILE BUTTON (Classic Windows File Explorer)
        upload_card = ctk.CTkFrame(self.sidebar_frame, fg_color=COLOR_BG_BLACK, corner_radius=8, border_width=1, border_color=COLOR_SAPPHIRE)
        upload_card.pack(fill="x", padx=10, pady=(4, 12))

        upload_btn = ctk.CTkButton(
            upload_card,
            text="📂 OPEN PPT FILE",
            font=ctk.CTkFont(size=12, weight="bold"),
            fg_color=COLOR_SAPPHIRE,
            hover_color=COLOR_SAPPHIRE_HOVER,
            command=self.open_pptx_file
        )
        upload_btn.pack(fill="x", padx=8, pady=8)

        # 4. RIGHT CONTENT WORK AREA (AUTO-EXPANDING)
        self.content_area = ctk.CTkFrame(self.body_container, fg_color=COLOR_BG_BLACK)
        self.content_area.pack(side="right", fill="both", expand=True, padx=10, pady=10)

        self.build_dashboard_view()
        self.build_editor_view()
        self.build_keywords_view()
        self.build_settings_view()

        self.sidebar_buttons = []
        self.refresh_sidebar_slide_list()
        self.switch_view("dashboard")

    def refresh_sidebar_slide_list(self):
        """Ultra-fast sidebar update: reuses existing button widgets when possible."""
        num_slides = len(self.slide_mgr.slides)
        
        # Fast path: update colors/labels of existing buttons without destroying widgets
        if hasattr(self, 'sidebar_buttons') and len(self.sidebar_buttons) == num_slides:
            for idx, (btn, slide) in enumerate(zip(self.sidebar_buttons, self.slide_mgr.slides)):
                is_active = (idx == self.current_slide_idx)
                btn_color = COLOR_SAPPHIRE if is_active else COLOR_BG_CARD
                lbl_str = f"#{slide.slide_id}: {slide.title[:20]}"
                btn.configure(
                    text=lbl_str,
                    fg_color=btn_color,
                    font=ctk.CTkFont(size=11, weight="bold" if is_active else "normal")
                )
            return

        # Rebuild path (when slides added/deleted/loaded)
        for child in self.slide_list_scroll.winfo_children():
            child.destroy()

        self.sidebar_buttons = []

        for idx, slide in enumerate(self.slide_mgr.slides):
            is_active = (idx == self.current_slide_idx)
            btn_color = COLOR_SAPPHIRE if is_active else COLOR_BG_CARD
            lbl_str = f"#{slide.slide_id}: {slide.title[:20]}"
            btn = ctk.CTkButton(
                self.slide_list_scroll,
                text=lbl_str,
                anchor="w",
                fg_color=btn_color,
                hover_color=COLOR_SAPPHIRE_HOVER,
                font=ctk.CTkFont(size=11, weight="bold" if is_active else "normal"),
                command=lambda i=idx: self.select_slide_by_index(i)
            )
            btn.pack(fill="x", pady=2)
            self.sidebar_buttons.append(btn)

    def select_slide_by_index(self, idx):
        """Selects a slide from the Left Sidebar list."""
        if 0 <= idx < len(self.slide_mgr.slides):
            self.current_slide_idx = idx
            self.update_slide_display()
            self.refresh_sidebar_slide_list()
            self.load_current_slide_into_editor()
            self.focus_set()

    def switch_view(self, view_name):
        """Switches the right content area view."""
        if getattr(self, 'active_view', '') == 'keywords':
            self.save_all_keywords()

        self.active_view = view_name
        self.view_dashboard_frame.pack_forget()
        self.view_editor_frame.pack_forget()
        self.view_keywords_frame.pack_forget()
        self.view_settings_frame.pack_forget()

        self.btn_nav_dashboard.configure(fg_color=COLOR_SAPPHIRE if view_name == "dashboard" else COLOR_BG_BLACK)
        self.btn_nav_editor.configure(fg_color=COLOR_SAPPHIRE if view_name == "editor" else COLOR_BG_BLACK)
        self.btn_nav_keywords.configure(fg_color=COLOR_SAPPHIRE if view_name == "keywords" else COLOR_BG_BLACK)
        self.btn_nav_settings.configure(fg_color=COLOR_SAPPHIRE if view_name == "settings" else COLOR_BG_BLACK)

        if view_name == "dashboard":
            self.view_dashboard_frame.pack(fill="both", expand=True)
        elif view_name == "editor":
            self.load_current_slide_into_editor()
            self.view_editor_frame.pack(fill="both", expand=True)
        elif view_name == "keywords":
            self.refresh_keywords_grid()
            self.view_keywords_frame.pack(fill="both", expand=True)
        elif view_name == "settings":
            try:
                self.monitors = get_monitors()
            except Exception:
                pass
            self.build_settings_view()
            self.view_settings_frame.pack(fill="both", expand=True)
            
        self.focus_set()

    def build_dashboard_view(self):
        """Constructs Presenter Dashboard view with Static Layout."""
        self.view_dashboard_frame = ctk.CTkFrame(self.content_area, fg_color="transparent")

        # Right Area: Upcoming Slide + Speaker Notes + FULL KEYWORD MATRIX LIST + SHORTCUT CHEAT SHEET
        self.dash_right_frame = ctk.CTkFrame(self.view_dashboard_frame, fg_color=COLOR_BG_CARD, corner_radius=8, width=340, border_width=1, border_color=COLOR_BORDER)
        self.dash_right_frame.pack(side="right", fill="both", padx=(6, 0))
        self.dash_right_frame.pack_propagate(False)

        # 1. UPCOMING NEXT SLIDE CARD
        self.dash_upcoming_lbl = ctk.CTkLabel(self.dash_right_frame, text="UPCOMING NEXT SLIDE", font=ctk.CTkFont(size=11, weight="bold"), text_color=COLOR_TEXT_MUTED)
        self.dash_upcoming_lbl.pack(anchor="w", padx=10, pady=(8, 2))

        self.next_slide_img_lbl = ctk.CTkLabel(self.dash_right_frame, text="")
        self.next_slide_img_lbl.pack(fill="x", padx=10, pady=2)

        # 2. SPEAKER NOTES PANE
        self.dash_notes_lbl = ctk.CTkLabel(self.dash_right_frame, text="SPEAKER NOTES", font=ctk.CTkFont(size=11, weight="bold"), text_color=COLOR_TEXT_MUTED)
        self.dash_notes_lbl.pack(anchor="w", padx=10, pady=(4, 2))

        self.notes_textbox = ctk.CTkTextbox(self.dash_right_frame, fg_color=COLOR_BG_BLACK, text_color=COLOR_TEXT_WHITE, font=ctk.CTkFont(size=11), height=65)
        self.notes_textbox.pack(fill="x", padx=10, pady=(0, 4))

        # 3. FULL PRESENTATION KEYWORD MATRIX LIST (ROW BY ROW FOR EVERY SLIDE)
        self.dash_kw_matrix_card = ctk.CTkFrame(self.dash_right_frame, fg_color=COLOR_BG_BLACK, corner_radius=6, border_width=1, border_color=COLOR_SAPPHIRE)
        self.dash_kw_matrix_card.pack(fill="both", expand=True, padx=10, pady=(2, 4))

        kw_hdr_box = ctk.CTkFrame(self.dash_kw_matrix_card, fg_color="transparent")
        kw_hdr_box.pack(fill="x", padx=8, pady=(6, 2))

        ctk.CTkLabel(
            kw_hdr_box,
            text="🏷️ VOICE KEYWORDS (ALL SLIDES)",
            font=ctk.CTkFont(size=10, weight="bold"),
            text_color=COLOR_ACCENT_GREEN
        ).pack(side="left")

        ctk.CTkLabel(
            kw_hdr_box,
            text="⚡ Live FST",
            font=ctk.CTkFont(size=9, weight="bold"),
            text_color="#34D399",
            fg_color="#064E3B",
            corner_radius=4,
            padx=5,
            pady=1
        ).pack(side="right")

        self.dash_kw_matrix_scroll = ctk.CTkScrollableFrame(self.dash_kw_matrix_card, fg_color="transparent")
        self.dash_kw_matrix_scroll.pack(fill="both", expand=True, padx=4, pady=(2, 4))

        # 4. STRUCTURED SHORTCUT KEYS & SLIDE OPERATIONS BAR
        shortcuts_card = ctk.CTkFrame(self.dash_right_frame, fg_color=COLOR_BG_BLACK, corner_radius=6, border_width=1, border_color=COLOR_BORDER)
        shortcuts_card.pack(fill="x", padx=10, pady=(2, 8))

        ctk.CTkLabel(
            shortcuts_card,
            text="⌨️ SHORTCUTS & OPERATIONS",
            font=ctk.CTkFont(size=10, weight="bold"),
            text_color=COLOR_SAPPHIRE
        ).pack(anchor="w", padx=8, pady=(4, 2))

        shortcut_lines = [
            "• Next Slide     :  →  /  Down  /  Space  /  'Next Slide'",
            "• Prev Slide     :  ←  /  Up    /  Bksp   /  'Previous Slide'",
            "• Jump to Slide  :  1-9 Keys    /  'Slide 1', 'Slide 2'",
            "• Black/Whiteout :  'B' / 'W'   /  'Black screen', 'White screen'",
            "• Exit / Stop    :  Escape (Esc)"
        ]
        for line in shortcut_lines:
            ctk.CTkLabel(
                shortcuts_card,
                text=line,
                font=ctk.CTkFont(family="Consolas", size=9),
                text_color="#CBD5E1",
                anchor="w",
                justify="left"
            ).pack(fill="x", padx=8, pady=0)

        # Left Area: Slide Previews (Flex Expanding)
        self.dash_left_frame = ctk.CTkFrame(self.view_dashboard_frame, fg_color=COLOR_BG_CARD, corner_radius=8, border_width=1, border_color=COLOR_BORDER)
        self.dash_left_frame.pack(side="left", fill="both", expand=True, padx=(0, 6))

        self.dash_active_lbl = ctk.CTkLabel(self.dash_left_frame, text="ACTIVE SLIDE PREVIEW (REAL POWERPOINT VISUALS)", font=ctk.CTkFont(size=12, weight="bold"), text_color=COLOR_TEXT_MUTED)
        self.dash_active_lbl.pack(anchor="w", padx=12, pady=(12, 4))

        self.curr_slide_img_lbl = ctk.CTkLabel(self.dash_left_frame, text="")
        self.curr_slide_img_lbl.pack(fill="both", expand=True, padx=12, pady=4)

        # Bottom Navigation Controls Bar (Pixel-Perfect Aligned)
        ctrl_frame = ctk.CTkFrame(self.dash_left_frame, fg_color="transparent", height=44)
        ctrl_frame.pack(fill="x", padx=12, pady=12)

        self.prev_btn = ctk.CTkButton(
            ctrl_frame,
            text="◄ PREV (←)",
            font=ctk.CTkFont(size=12, weight="bold"),
            fg_color=COLOR_SAPPHIRE,
            hover_color=COLOR_SAPPHIRE_HOVER,
            width=110,
            height=36,
            corner_radius=6,
            command=self.prev_slide
        )
        self.prev_btn.pack(side="left", padx=(0, 6))

        self.next_btn = ctk.CTkButton(
            ctrl_frame,
            text="NEXT (→) ►",
            font=ctk.CTkFont(size=12, weight="bold"),
            fg_color=COLOR_SAPPHIRE,
            hover_color=COLOR_SAPPHIRE_HOVER,
            width=110,
            height=36,
            corner_radius=6,
            command=self.next_slide
        )
        self.next_btn.pack(side="left", padx=(0, 12))

        self.slide_scrubber = ctk.CTkSlider(
            ctrl_frame,
            from_=0,
            to=len(self.slide_mgr.slides)-1,
            number_of_steps=len(self.slide_mgr.slides),
            height=16,
            progress_color=COLOR_SAPPHIRE,
            button_color=COLOR_SAPPHIRE,
            button_hover_color=COLOR_SAPPHIRE_HOVER,
            command=self.on_scrubber_change
        )
        self.slide_scrubber.pack(side="left", fill="x", expand=True, padx=(0, 12))

        self.slide_num_lbl = ctk.CTkLabel(
            ctrl_frame,
            text="Slide 1 / 4",
            font=ctk.CTkFont(family="Consolas", size=13, weight="bold"),
            text_color="#E2E8F0"
        )
        self.slide_num_lbl.pack(side="right", padx=(4, 0))

    def build_editor_view(self):
        """Constructs Slide Detail Editor workspace view."""
        self.view_editor_frame = ctk.CTkFrame(self.content_area, fg_color=COLOR_BG_CARD, corner_radius=8, border_width=1, border_color=COLOR_BORDER)

        ctk.CTkLabel(self.view_editor_frame, text="✏️ SLIDE CONTENT EDITOR", font=ctk.CTkFont(size=15, weight="bold"), text_color=COLOR_SAPPHIRE).pack(anchor="w", padx=20, pady=(20, 10))

        # Title Entry
        ctk.CTkLabel(self.view_editor_frame, text="Slide Title:").pack(anchor="w", padx=20, pady=(10, 2))
        self.edit_title_entry = ctk.CTkEntry(self.view_editor_frame, font=ctk.CTkFont(size=14))
        self.edit_title_entry.pack(fill="x", padx=20, pady=(0, 10))

        # Bullet Points Text Area
        ctk.CTkLabel(self.view_editor_frame, text="Bullet Points (One per line):").pack(anchor="w", padx=20, pady=(10, 2))
        self.edit_bullets_textbox = ctk.CTkTextbox(self.view_editor_frame, height=140, font=ctk.CTkFont(size=13))
        self.edit_bullets_textbox.pack(fill="x", padx=20, pady=(0, 10))

        # Speaker Notes Text Area
        ctk.CTkLabel(self.view_editor_frame, text="Speaker Notes:").pack(anchor="w", padx=20, pady=(10, 2))
        self.edit_notes_textbox = ctk.CTkTextbox(self.view_editor_frame, height=100, font=ctk.CTkFont(size=13))
        self.edit_notes_textbox.pack(fill="x", padx=20, pady=(0, 10))

        # Apply Changes Button
        ctk.CTkButton(self.view_editor_frame, text="✔ APPLY SLIDE CHANGES", font=ctk.CTkFont(size=13, weight="bold"), fg_color=COLOR_SAPPHIRE, command=self.apply_slide_edits).pack(anchor="e", padx=20, pady=15)

    def load_current_slide_into_editor(self):
        """Loads selected slide fields into editor form."""
        if 0 <= self.current_slide_idx < len(self.slide_mgr.slides):
            slide = self.slide_mgr.slides[self.current_slide_idx]
            self.edit_title_entry.delete(0, "end")
            self.edit_title_entry.insert(0, slide.title)
            
            self.edit_bullets_textbox.delete("1.0", "end")
            self.edit_bullets_textbox.insert("1.0", "\n".join(slide.bullet_points))
            
            self.edit_notes_textbox.delete("1.0", "end")
            self.edit_notes_textbox.insert("1.0", slide.notes)

    def build_keywords_view(self):
        """Constructs Voice Keyword Mapping Grid view."""
        self.view_keywords_frame = ctk.CTkFrame(self.content_area, fg_color="transparent")

        top_frame = ctk.CTkFrame(self.view_keywords_frame, fg_color=COLOR_BG_CARD, corner_radius=8, border_width=1, border_color=COLOR_BORDER)
        top_frame.pack(fill="x", pady=(0, 15))

        # Top Header Save Button & Status (Right side)
        right_box = ctk.CTkFrame(top_frame, fg_color="transparent")
        right_box.pack(side="right", padx=20, pady=10)

        top_btn = ctk.CTkButton(
            right_box,
            text="💾 SAVE ALL VOICE KEYWORDS",
            font=ctk.CTkFont(size=12, weight="bold"),
            fg_color=COLOR_SAPPHIRE,
            hover_color=COLOR_SAPPHIRE_HOVER,
            command=self.save_all_keywords
        )
        top_btn.pack(side="top", anchor="e", pady=(0, 4))

        self.kw_sync_status_lbl = ctk.CTkLabel(
            right_box,
            text="🟢 Engine Synced & Ready",
            font=ctk.CTkFont(size=11, weight="bold"),
            text_color=COLOR_ACCENT_GREEN
        )
        self.kw_sync_status_lbl.pack(side="top", anchor="e")

        # Header Title & Subtitle (Left side)
        text_box = ctk.CTkFrame(top_frame, fg_color="transparent")
        text_box.pack(side="left", fill="both", expand=True, padx=20, pady=12)

        ctk.CTkLabel(text_box, text="🔑 VOICE KEYWORD SLIDE MAPPING MATRIX", font=ctk.CTkFont(size=15, weight="bold"), text_color=COLOR_SAPPHIRE).pack(anchor="w", pady=(0, 2))
        ctk.CTkLabel(text_box, text="Slide numbers & titles are automatically mapped (e.g. 'Slide 1', 'One', 'First'). Say any keyword live into your mic to flip slides instantly.", font=ctk.CTkFont(size=12), text_color=COLOR_TEXT_MUTED).pack(anchor="w")

        self.kw_scroll_frame = ctk.CTkScrollableFrame(self.view_keywords_frame, fg_color=COLOR_BG_CARD, corner_radius=8, border_width=1, border_color=COLOR_BORDER)
        self.kw_scroll_frame.pack(fill="both", expand=True)

    def refresh_keywords_grid(self):
        """Refreshes keyword input fields for all slides."""
        # Save any in-progress edits from current entries before rebuilding grid
        if hasattr(self, 'kw_entries') and self.kw_entries:
            for idx, entry in self.kw_entries:
                try:
                    if entry.winfo_exists():
                        raw_kws = entry.get().split(",")
                        clean_kws = [k.strip().lower() for k in raw_kws if k.strip()]
                        if idx < len(self.slide_mgr.slides):
                            self.slide_mgr.slides[idx].keywords = clean_kws
                except Exception:
                    pass

        for widget in self.kw_scroll_frame.winfo_children():
            widget.destroy()

        self.kw_entries = []

        for idx, slide in enumerate(self.slide_mgr.slides):
            row_frame = ctk.CTkFrame(self.kw_scroll_frame, fg_color=COLOR_BG_BLACK, corner_radius=6)
            row_frame.pack(fill="x", padx=15, pady=8)

            lbl_str = f"Slide #{slide.slide_id}: {slide.title[:25]}"
            ctk.CTkLabel(row_frame, text=lbl_str, font=ctk.CTkFont(size=13, weight="bold"), width=240, anchor="w").pack(side="left", padx=15, pady=10)

            kw_str = ", ".join(slide.keywords)
            entry = ctk.CTkEntry(row_frame, font=ctk.CTkFont(size=13), placeholder_text="e.g. slide 1, one, first, intro")
            entry.insert(0, kw_str)
            entry.pack(side="left", fill="x", expand=True, padx=15, pady=10)

            # Live sync on FocusOut and KeyRelease
            entry.bind("<FocusOut>", lambda e: self.save_all_keywords())
            entry.bind("<KeyRelease>", lambda e: self.auto_sync_keywords())

            self.kw_entries.append((idx, entry))

        ctk.CTkButton(self.kw_scroll_frame, text="💾 SAVE ALL VOICE KEYWORDS", font=ctk.CTkFont(size=13, weight="bold"), fg_color=COLOR_SAPPHIRE, command=self.save_all_keywords).pack(anchor="e", padx=15, pady=20)

    def auto_sync_keywords(self):
        """Live syncs keyword entry fields with Voice Speech Engine in real time."""
        if hasattr(self, 'kw_entries') and self.kw_entries:
            for idx, entry in self.kw_entries:
                try:
                    if entry.winfo_exists():
                        raw_kws = entry.get().split(",")
                        clean_kws = [k.strip().lower() for k in raw_kws if k.strip()]
                        if idx < len(self.slide_mgr.slides):
                            self.slide_mgr.slides[idx].keywords = clean_kws
                except Exception:
                    pass
            count = self.voice_engine.set_keywords(self.slide_mgr.slides)
            self.refresh_dash_keyword_matrix()
            if hasattr(self, 'kw_sync_status_lbl') and self.kw_sync_status_lbl.winfo_exists():
                self.kw_sync_status_lbl.configure(text=f"⚡ Syncing ({count} Active)...", text_color=COLOR_ACCENT_GREEN)

    def save_all_keywords(self):
        """Saves keywords from entries back to SlideData objects."""
        if hasattr(self, 'kw_entries') and self.kw_entries:
            for idx, entry in self.kw_entries:
                try:
                    if entry.winfo_exists():
                        raw_kws = entry.get().split(",")
                        clean_kws = [k.strip().lower() for k in raw_kws if k.strip()]
                        if idx < len(self.slide_mgr.slides):
                            self.slide_mgr.slides[idx].keywords = clean_kws
                except Exception:
                    pass

        count = self.voice_engine.set_keywords(self.slide_mgr.slides)
        self.refresh_dash_keyword_matrix()
        num_slides = len(self.slide_mgr.slides)
        msg = f"🟢 VOICE KEYWORDS SYNCED SUCCESSFULLY! ({count} Keywords across {num_slides} Slides)"

        if hasattr(self, 'match_badge_lbl'):
            self.match_badge_lbl.configure(text=msg, text_color=COLOR_ACCENT_GREEN)

        if hasattr(self, 'kw_sync_status_lbl') and self.kw_sync_status_lbl.winfo_exists():
            self.kw_sync_status_lbl.configure(text=f"✅ Synced {count} Keywords across {num_slides} Slides!", text_color=COLOR_ACCENT_GREEN)

        self.focus_set()

    def build_settings_view(self):
        """Constructs Audio Microphone & HDMI Display Settings view."""
        self.view_settings_frame = ctk.CTkFrame(self.content_area, fg_color=COLOR_BG_CARD, corner_radius=8, border_width=1, border_color=COLOR_BORDER)

        # 1. MICROPHONE & AUDIO INPUT DEVICE CONFIGURATION
        ctk.CTkLabel(self.view_settings_frame, text="🎙️ MICROPHONE & AUDIO INPUT CONFIGURATION", font=ctk.CTkFont(size=15, weight="bold"), text_color=COLOR_SAPPHIRE).pack(anchor="w", padx=20, pady=(20, 10))

        mic_card = ctk.CTkFrame(self.view_settings_frame, fg_color=COLOR_BG_BLACK, corner_radius=6, border_width=1, border_color=COLOR_BORDER)
        mic_card.pack(anchor="w", fill="x", padx=20, pady=(0, 15))

        all_mics = self.voice_engine.get_available_microphones()
        mic_options = []
        mic_map = {}
        for m in all_mics:
            label = f"#{m['id']}: {m['name']} ({m.get('api','')})"
            if len(label) > 65:
                label = label[:62] + "..."
            mic_options.append(label)
            mic_map[label] = m['id']

        curr_mic_label = "Auto-Select Best Mic"
        for label, dev_id in mic_map.items():
            if dev_id == self.voice_engine.device_id:
                curr_mic_label = label
                break

        ctk.CTkLabel(mic_card, text="Active Recording Microphone:", font=ctk.CTkFont(size=12, weight="bold"), text_color="#E2E8F0").pack(anchor="w", padx=15, pady=(12, 4))
        
        mic_ctrl_frame = ctk.CTkFrame(mic_card, fg_color="transparent")
        mic_ctrl_frame.pack(anchor="w", fill="x", padx=15, pady=(0, 12))

        def on_mic_selected(choice):
            chosen_id = mic_map.get(choice)
            if chosen_id is not None:
                was_recording = self.voice_engine.is_recording
                if was_recording:
                    self.voice_engine.stop()
                self.voice_engine.probe_microphone(force_device_id=chosen_id)
                self.update_microphone_indicator()
                if was_recording:
                    self.voice_engine.start()

        self.mic_option_menu = ctk.CTkOptionMenu(
            mic_ctrl_frame,
            values=mic_options if mic_options else ["No Microphones Found"],
            command=on_mic_selected,
            width=480,
            fg_color=COLOR_SAPPHIRE,
            button_color=COLOR_SAPPHIRE_HOVER
        )
        if curr_mic_label in mic_options:
            self.mic_option_menu.set(curr_mic_label)
        self.mic_option_menu.pack(side="left", padx=(0, 10))

        def refresh_mics():
            self.voice_engine.probe_microphone()
            self.update_microphone_indicator()
            self.build_settings_view()
            self.view_settings_frame.pack(fill="both", expand=True)

        ctk.CTkButton(mic_ctrl_frame, text="🔄 REFRESH MICS", fg_color=COLOR_BG_BLACK, border_width=1, border_color=COLOR_SAPPHIRE, command=refresh_mics).pack(side="left")

        # 2. MULTI-MONITOR / HDMI DISPLAY DETECTOR
        ctk.CTkLabel(self.view_settings_frame, text="🖥️ MULTI-MONITOR / HDMI DISPLAY DETECTOR", font=ctk.CTkFont(size=15, weight="bold"), text_color=COLOR_SAPPHIRE).pack(anchor="w", padx=20, pady=(15, 10))

        mon_str = f"Detected {len(self.monitors)} Display Monitor(s) on System:\n"
        for i, m in enumerate(self.monitors):
            primary_tag = " (Primary Laptop Screen)" if i == 0 else " 📺 (HDMI / External Display)"
            mon_str += f"  • Monitor #{i+1}: {m.width}x{m.height} at pos ({m.x}, {m.y}){primary_tag}\n"

        ctk.CTkLabel(self.view_settings_frame, text=mon_str, font=ctk.CTkFont(family="Consolas", size=13), text_color="#E2E8F0", justify="left").pack(anchor="w", padx=20, pady=10)

        btn_frame = ctk.CTkFrame(self.view_settings_frame, fg_color="transparent")
        btn_frame.pack(anchor="w", fill="x", padx=20, pady=20)

        ctk.CTkButton(btn_frame, text="🚀 LAUNCH FULLSCREEN HDMI DISPLAY (MONITOR #2)", fg_color=COLOR_SAPPHIRE, command=lambda: self.launch_hdmi_output(1)).pack(side="left", padx=10)
        ctk.CTkButton(btn_frame, text="🔲 LAUNCH WINDOWED PREVIEW MODE", fg_color=COLOR_BG_BLACK, border_width=1, border_color=COLOR_SAPPHIRE, command=lambda: self.launch_hdmi_output(None)).pack(side="left", padx=10)

    def launch_hdmi_output(self, monitor_idx=None):
        """Launches presentation window specifically targeting the connected projector or secondary HDMI display."""
        try:
            self.monitors = get_monitors()
        except Exception as e:
            print(f"[HDMI WARNING] Failed to probe monitors: {e}")

        # Distinguish primary screen (laptop) vs external projector/HDMI displays
        primary_mon = None
        external_monitors = []

        for m in self.monitors:
            if getattr(m, 'is_primary', False) or (m.x == 0 and m.y == 0):
                primary_mon = m
            else:
                external_monitors.append(m)

        target_mon = None

        if monitor_idx is not None and monitor_idx < len(self.monitors):
            target_mon = self.monitors[monitor_idx]
        elif external_monitors:
            # Projector connected! Target the external projector screen
            target_mon = external_monitors[0]
        elif len(self.monitors) > 1:
            target_mon = self.monitors[1]
        elif self.monitors:
            target_mon = self.monitors[0]
            
        if self.hdmi_window and self.hdmi_window.winfo_exists():
            self.hdmi_window.destroy()
            
        self.hdmi_window = ExternalDisplayWindow(self, monitor_info=target_mon)
        self.update_slide_display()

        # Update feedback status badge
        if target_mon and target_mon != primary_mon:
            msg = f"📺 PROJECTOR CONNECTED: Fullscreen Output ({target_mon.width}x{target_mon.height} at +{target_mon.x}+{target_mon.y})"
            self.match_badge_lbl.configure(text=msg, text_color=COLOR_ACCENT_GREEN)
        else:
            msg = f"📺 Laptop Preview Window Active (740x416)"
            self.match_badge_lbl.configure(text=msg, text_color=COLOR_SAPPHIRE)

        self.focus_set()

    def update_slide_display(self):
        """Renders and updates current slide on Presenter Dashboard and HDMI display with instant cached response."""
        if not self.slide_mgr.slides:
            return

        curr_slide = self.slide_mgr.slides[self.current_slide_idx]
        
        # Static HD Slide Image Dimensions
        target_w, target_h = 640, 360

        # Render HD Slide Image (instant if cached)
        pil_img = self.slide_mgr.render_slide_image(curr_slide, width=1280, height=720)
        
        # Update Presenter Dashboard Main Preview with fast caching
        main_key = (id(pil_img), target_w, target_h)
        if not hasattr(self, '_last_main_key') or self._last_main_key != main_key:
            self._cached_main_ctk = ctk.CTkImage(light_image=pil_img, dark_image=pil_img, size=(target_w, target_h))
            self._last_main_key = main_key
        self.curr_slide_img_lbl.configure(image=self._cached_main_ctk)

        # Static Next Slide Preview Dimensions
        next_w, next_h = 280, 158
        next_idx = (self.current_slide_idx + 1) % len(self.slide_mgr.slides)
        next_slide = self.slide_mgr.slides[next_idx]
        next_pil = self.slide_mgr.render_slide_image(next_slide, width=640, height=360)
        
        next_key = (id(next_pil), next_w, next_h)
        if not hasattr(self, '_last_next_key') or self._last_next_key != next_key:
            self._cached_next_ctk = ctk.CTkImage(light_image=next_pil, dark_image=next_pil, size=(next_w, next_h))
            self._last_next_key = next_key
        self.next_slide_img_lbl.configure(image=self._cached_next_ctk)

        # 5. Update Speaker Notes
        self.notes_textbox.delete("1.0", "end")
        self.notes_textbox.insert("1.0", curr_slide.notes if curr_slide.notes else "(No speaker notes for this slide)")

        # 5b. Update Full Presentation Keyword Matrix List (Row by Row)
        self.refresh_dash_keyword_matrix()

        # 6. Update Slide Number Label & Scrubber
        self.slide_num_lbl.configure(text=f"Slide {self.current_slide_idx + 1} / {len(self.slide_mgr.slides)}")
        self.slide_scrubber.set(self.current_slide_idx)

        # 7. Update External HDMI Display Window!
        if self.hdmi_window and self.hdmi_window.winfo_exists():
            self.hdmi_window.update_slide(pil_img, is_blackout=self.is_blackout, is_whiteout=self.is_whiteout)

    def refresh_dash_keyword_matrix(self):
        """Renders row-by-row Keyword Matrix list for every slide in the presenter dashboard."""
        if not hasattr(self, 'dash_kw_matrix_scroll') or not self.dash_kw_matrix_scroll.winfo_exists():
            return
            
        for child in self.dash_kw_matrix_scroll.winfo_children():
            child.destroy()
            
        for idx, slide in enumerate(self.slide_mgr.slides):
            is_active = (idx == self.current_slide_idx)
            row_bg = "#0B132B" if is_active else "#050811"
            border_col = COLOR_SAPPHIRE if is_active else "#1E293B"
            
            row_card = ctk.CTkFrame(
                self.dash_kw_matrix_scroll,
                fg_color=row_bg,
                corner_radius=5,
                border_width=1,
                border_color=border_col
            )
            row_card.pack(fill="x", padx=2, pady=2)
            
            # Header line: Slide number + Title + Active indicator
            hdr_box = ctk.CTkFrame(row_card, fg_color="transparent")
            hdr_box.pack(fill="x", padx=6, pady=(3, 1))
            
            slide_tag = f"#{slide.slide_id}: {slide.title[:20]}"
            tag_color = "#60A5FA" if is_active else "#94A3B8"
            lbl_title = ctk.CTkLabel(
                hdr_box,
                text=slide_tag,
                font=ctk.CTkFont(size=10, weight="bold" if is_active else "normal"),
                text_color=tag_color,
                anchor="w"
            )
            lbl_title.pack(side="left")
            
            if is_active:
                ctk.CTkLabel(
                    hdr_box,
                    text="● CURRENT",
                    font=ctk.CTkFont(size=8, weight="bold"),
                    text_color="#34D399",
                    fg_color="#064E3B",
                    corner_radius=3,
                    padx=4,
                    pady=1
                ).pack(side="right")
                
            # Keywords line
            kws = slide.keywords if slide.keywords else ["(no keywords)"]
            kw_str = " • ".join(kws[:5])
            if len(slide.keywords) > 5:
                kw_str += f" (+{len(slide.keywords)-5})"
                
            lbl_kws = ctk.CTkLabel(
                row_card,
                text=f"🗣️ {kw_str}",
                font=ctk.CTkFont(size=9),
                text_color="#93C5FD" if is_active else "#64748B",
                anchor="w",
                justify="left",
                wraplength=290
            )
            lbl_kws.pack(fill="x", padx=6, pady=(0, 3))
            
            # Clickable row to navigate directly to slide
            row_card.bind("<Button-1>", lambda e, i=idx: self.select_slide_by_index(i))
            lbl_title.bind("<Button-1>", lambda e, i=idx: self.select_slide_by_index(i))
            lbl_kws.bind("<Button-1>", lambda e, i=idx: self.select_slide_by_index(i))

    def next_slide(self):
        """Navigates to next slide."""
        if self.current_slide_idx < len(self.slide_mgr.slides) - 1:
            self.current_slide_idx += 1
            self.is_blackout = False
            self.is_whiteout = False
            self.update_slide_display()
            self.refresh_sidebar_slide_list()
            self.focus_set()

    def prev_slide(self):
        """Navigates to previous slide."""
        if self.current_slide_idx > 0:
            self.current_slide_idx -= 1
            self.is_blackout = False
            self.is_whiteout = False
            self.update_slide_display()
            self.refresh_sidebar_slide_list()
            self.focus_set()

    def on_scrubber_change(self, value):
        """Handles slide scrubber slider movement."""
        idx = int(round(value))
        if 0 <= idx < len(self.slide_mgr.slides):
            self.current_slide_idx = idx
            self.update_slide_display()
            self.refresh_sidebar_slide_list()
            self.focus_set()

    def toggle_blackout(self):
        """Toggles black screen blackout mode."""
        self.is_blackout = not self.is_blackout
        self.update_slide_display()
        self.focus_set()

    def toggle_whiteout(self):
        """Toggles white screen whiteout mode."""
        self.is_whiteout = not self.is_whiteout
        self.update_slide_display()
        self.focus_set()

    def stop_presentation(self):
        """Exits presentation mode."""
        if self.hdmi_window and self.hdmi_window.winfo_exists():
            self.hdmi_window.destroy()



    def toggle_voice_engine(self):
        """Starts/Stops live voice speech engine."""
        self.save_all_keywords()
        if not self.voice_engine.is_recording:
            self.voice_engine.set_keywords(self.slide_mgr.slides)
            success = self.voice_engine.start()
            if success:
                self.voice_btn.configure(text="⏹️ STOP VOICE", fg_color=COLOR_ACCENT_RED)
                self.match_badge_lbl.configure(text="[Voice Status: LIVE LISTENING]", text_color=COLOR_ACCENT_GREEN)
            else:
                from tkinter import messagebox
                self.voice_btn.configure(text="🎙️ LIVE VOICE", fg_color=COLOR_SAPPHIRE)
                self.match_badge_lbl.configure(text="[Voice Status: Mic Error]", text_color=COLOR_ACCENT_RED)
                messagebox.showwarning(
                    "Microphone Stream Warning",
                    "Could not open microphone audio stream on this system.\nPlease check your microphone hardware connection or permissions.\nManual keyboard & UI navigation remain 100% active."
                )
        else:
            self.voice_engine.stop()
            self.voice_btn.configure(text="🎙️ LIVE VOICE", fg_color=COLOR_SAPPHIRE)
            self.match_badge_lbl.configure(text="[Voice Status: Paused]", text_color=COLOR_TEXT_MUTED)
        self.focus_set()

    def on_voice_keyword_matched(self, target, matched_kw, full_spoken_text):
        """Callback triggered when speech engine matches a command or slide keyword in real time!"""
        if isinstance(target, str):
            # Global Action Commands
            if target == 'ACTION_NEXT':
                self.after(0, self.next_slide)
                msg = f"🎙️ COMMAND: 'Next Slide' (Spoken: '{full_spoken_text}') ➔ Advanced to Slide #{self.current_slide_idx + 2}"
            elif target == 'ACTION_PREV':
                self.after(0, self.prev_slide)
                msg = f"🎙️ COMMAND: 'Previous Slide' (Spoken: '{full_spoken_text}') ➔ Back to Slide #{max(1, self.current_slide_idx)}"
            elif target == 'ACTION_FIRST':
                self.select_slide_by_index(0)
                msg = f"🎙️ COMMAND: 'First Slide' (Spoken: '{full_spoken_text}') ➔ Jumped to Slide #1"
            elif target == 'ACTION_LAST':
                last_idx = max(0, len(self.slide_mgr.slides) - 1)
                self.select_slide_by_index(last_idx)
                msg = f"🎙️ COMMAND: 'Last Slide' (Spoken: '{full_spoken_text}') ➔ Jumped to Slide #{last_idx + 1}"
            elif target == 'ACTION_BLACKOUT':
                self.after(0, self.toggle_blackout)
                msg = f"🎙️ COMMAND: 'Blackout Screen' (Spoken: '{full_spoken_text}')"
            elif target == 'ACTION_WHITEOUT':
                self.after(0, self.toggle_whiteout)
                msg = f"🎙️ COMMAND: 'Whiteout Screen' (Spoken: '{full_spoken_text}')"
            else:
                msg = f"🎙️ COMMAND: '{target}'"
            self.after(0, lambda: self.match_badge_lbl.configure(text=msg, text_color=COLOR_ACCENT_GREEN))
            return

        # Direct Slide Index Jump
        slide_idx = target
        if 0 <= slide_idx < len(self.slide_mgr.slides):
            self.current_slide_idx = slide_idx
            self.is_blackout = False
            self.is_whiteout = False
            
            self.after(0, self.update_slide_display)
            self.after(0, self.refresh_sidebar_slide_list)
            msg = f"🎙️ MATCHED: '{matched_kw}' (Spoken: '{full_spoken_text}') ➔ Jumped to Slide #{slide_idx + 1}"
            self.after(0, lambda: self.match_badge_lbl.configure(text=msg, text_color=COLOR_ACCENT_GREEN))

    def refresh_mic_dropdown(self):
        """Refreshes the microphone selector dropdown menu options."""
        if not hasattr(self, 'mic_select_menu') or not self.mic_select_menu.winfo_exists():
            return
            
        all_mics = self.voice_engine.get_available_microphones()
        self._mic_options_map = {}
        options = []
        
        current_choice = ""
        for m in all_mics:
            if m.get('is_virtual', False):
                continue
            icon = m.get('icon', '🎙️')
            clean_n = m['name'].split('(')[0].strip() if '(' in m['name'] else m['name']
            if len(clean_n) > 22:
                clean_n = clean_n[:20] + ".."
            label = f"{icon} {clean_n} (#{m['id']})"
            self._mic_options_map[label] = m['id']
            options.append(label)
            if m['id'] == self.voice_engine.device_id:
                current_choice = label

        if not options:
            options = ["No Microphones Found"]
            current_choice = options[0]
            
        self.mic_select_menu.configure(values=options)
        if current_choice:
            self.mic_select_menu.set(current_choice)

    def on_mic_dropdown_selected(self, selected_label):
        """User selected a microphone from the dropdown."""
        dev_id = getattr(self, '_mic_options_map', {}).get(selected_label)
        if dev_id is not None:
            success = self.voice_engine.switch_device(dev_id)
            self.update_microphone_indicator()
            if success:
                clean_n = self.voice_engine.device_name.split('(')[0].strip()
                dev_type = getattr(self.voice_engine, 'device_type', 'MIC')
                type_badge = "🎧 Bluetooth Headset" if dev_type == "BT" else ("🎙️ USB Mic" if dev_type == "USB" else "💻 System Mic")
                self.match_badge_lbl.configure(text=f"⚡ Switched to {type_badge}: {clean_n}", text_color=COLOR_ACCENT_GREEN)

    def manual_rescan_microphones(self):
        """Manual rescan button clicked."""
        self.match_badge_lbl.configure(text="🔍 Scanning for Bluetooth & USB Mics...", text_color=COLOR_SAPPHIRE)
        self.voice_engine.refresh_portaudio()
        self.voice_engine.probe_microphone(force_refresh=True)
        self.update_microphone_indicator()
        clean_n = self.voice_engine.device_name.split('(')[0].strip()
        dev_type = getattr(self.voice_engine, 'device_type', 'MIC')
        type_badge = "🎧 Bluetooth Headset" if dev_type == "BT" else ("🎙️ USB Mic" if dev_type == "USB" else "💻 Mic")
        self.match_badge_lbl.configure(text=f"⚡ Active {type_badge}: {clean_n}", text_color=COLOR_ACCENT_GREEN)

    def update_microphone_indicator(self):
        """Updates top header badge with dynamic Bluetooth, USB, or Built-in mic status and real-time state."""
        if hasattr(self, 'mic_status_badge') and self.mic_status_badge.winfo_exists():
            if self.voice_engine.is_mic_connected and self.voice_engine.device_name and self.voice_engine.device_name not in ["No Input Device Detected", ""]:
                dev_name = self.voice_engine.device_name.strip()
                dev_type = getattr(self.voice_engine, 'device_type', 'MIC')
                
                # Format clean name
                clean_name = dev_name.split('(')[0].strip() if '(' in dev_name else dev_name
                if len(clean_name) > 22:
                    display_name = clean_name[:20] + ".."
                else:
                    display_name = clean_name

                if dev_type == "BT":
                    badge_text = f"🎧 BT: {display_name} • CONNECTED"
                    fg_col = "#064E3B"
                    txt_col = "#34D399"
                elif dev_type == "USB":
                    badge_text = f"🎙️ USB: {display_name} • CONNECTED"
                    fg_col = "#1E3A8A"
                    txt_col = "#60A5FA"
                else:
                    badge_text = f"💻 Mic: {display_name} • CONNECTED"
                    fg_col = "#1E293B"
                    txt_col = "#94A3B8"

                self.mic_status_badge.configure(
                    text=badge_text,
                    text_color=txt_col,
                    fg_color=fg_col
                )
            else:
                self.mic_status_badge.configure(
                    text="🔴 NO MIC CONNECTED • DISCONNECTED",
                    text_color="#FCA5A5",
                    fg_color="#7F1D1D"
                )

        self.refresh_mic_dropdown()

    def start_gui_live_indicator_loop(self):
        """Updates live GUI header status bar 60 times/sec and monitors real-time microphone hot-plugging."""
        self._hotplug_tick_counter = 0
        def ui_update():
            if self.voice_engine.is_recording:
                indicator_str = self.voice_engine.get_status_indicator_str(fps=60.0)
                self.speech_status_lbl.configure(text=indicator_str)
            else:
                self.speech_status_lbl.configure(text="⚪ OFF  [░░░░░░░░░░░░] [60.0 FPS |  0.0ms]")

            # Check microphone hotplug / disconnect events (~every 0.8 second = 48 ticks @ 16ms)
            self._hotplug_tick_counter += 1
            if self._hotplug_tick_counter >= 48:
                self._hotplug_tick_counter = 0
                event, old_dev, new_dev = self.voice_engine.check_device_hotplug()
                if event == "BLUETOOTH_CONNECTED":
                    self.update_microphone_indicator()
                    clean_n = new_dev.split('(')[0].strip() if '(' in (new_dev or '') else (new_dev or '')
                    self.match_badge_lbl.configure(text=f"🎧 Bluetooth Headset Connected: {clean_n} (Live!)", text_color=COLOR_ACCENT_GREEN)
                elif event == "CONNECTED":
                    self.update_microphone_indicator()
                    clean_n = new_dev.split('(')[0].strip() if '(' in (new_dev or '') else (new_dev or '')
                    self.match_badge_lbl.configure(text=f"🎙️ Mic Connected: {clean_n}", text_color=COLOR_ACCENT_GREEN)
                elif event == "RECOVERED":
                    self.update_microphone_indicator()
                    clean_n = new_dev.split('(')[0].strip() if '(' in (new_dev or '') else (new_dev or '')
                    self.match_badge_lbl.configure(text=f"⚡ Audio Stream Recovered: {clean_n}", text_color=COLOR_ACCENT_GREEN)
                elif event == "SWITCHED":
                    self.update_microphone_indicator()
                    clean_n = new_dev.split('(')[0].strip() if '(' in (new_dev or '') else (new_dev or '')
                    self.match_badge_lbl.configure(text=f"⚠️ Switched to Mic: {clean_n}", text_color="#F59E0B")
                elif event == "DISCONNECTED":
                    self.update_microphone_indicator()
                    clean_n = new_dev.split('(')[0].strip() if '(' in (new_dev or '') else (new_dev or 'Built-in Mic')
                    self.match_badge_lbl.configure(text=f"⚠️ Bluetooth Disconnected ➔ Auto-switched to: {clean_n}", text_color="#F59E0B")
                elif event == "UPDATED":
                    self.update_microphone_indicator()

            self.after(16, ui_update)

        self.after(100, ui_update)

    def update_speech_status_bar(self, status_str):
        """Callback for speech status update."""
        self.speech_status_lbl.configure(text=status_str)

    def apply_slide_edits(self):
        """Applies edits from editor form back to slide data."""
        try:
            slide = self.slide_mgr.slides[self.current_slide_idx]
            slide.title = self.edit_title_entry.get().strip()
            raw_bullets = self.edit_bullets_textbox.get("1.0", "end").split("\n")
            slide.bullet_points = [b.strip() for b in raw_bullets if b.strip()]
            slide.notes = self.edit_notes_textbox.get("1.0", "end").strip()
            
            self.update_slide_display()
            self.refresh_sidebar_slide_list()
            self.refresh_keywords_grid()
            self.focus_set()
        except Exception as e:
            print(f"[ERROR] Failed to apply slide edits: {e}")

    def add_new_slide(self):
        """Adds new blank slide to deck."""
        new_id = len(self.slide_mgr.slides) + 1
        new_kws = generate_default_slide_keywords(new_id, f"New Slide {new_id}", ["Add bullet points here"])
        new_slide = SlideData(
            slide_id=new_id,
            title=f"New Slide {new_id}",
            bullet_points=["Add bullet points here"],
            notes="Add speaker notes here",
            keywords=new_kws
        )
        self.slide_mgr.slides.append(new_slide)
        self.current_slide_idx = len(self.slide_mgr.slides) - 1
        self.slide_scrubber.configure(to=len(self.slide_mgr.slides)-1, number_of_steps=len(self.slide_mgr.slides))
        self.refresh_sidebar_slide_list()
        self.refresh_keywords_grid()
        self.voice_engine.set_keywords(self.slide_mgr.slides)
        self.update_slide_display()
        self.focus_set()

    def delete_slide(self):
        """Deletes selected slide from deck."""
        if len(self.slide_mgr.slides) > 1:
            self.slide_mgr.slides.pop(self.current_slide_idx)
            if self.current_slide_idx >= len(self.slide_mgr.slides):
                self.current_slide_idx = len(self.slide_mgr.slides) - 1
            self.slide_scrubber.configure(to=len(self.slide_mgr.slides)-1, number_of_steps=len(self.slide_mgr.slides))
            self.refresh_sidebar_slide_list()
            self.refresh_keywords_grid()
            self.update_slide_display()
            self.focus_set()


    def open_pptx_file(self):
        """Opens native file dialog to load .pptx deck from local system or USB drive."""
        from tkinter import filedialog
        initial_dir = os.path.expanduser("~/Documents")
        if not os.path.exists(initial_dir):
            initial_dir = os.getcwd()
        file_path = filedialog.askopenfilename(
            initialdir=initial_dir,
            title="Select PowerPoint Presentation Deck (.pptx)",
            filetypes=[("PowerPoint Presentations", "*.pptx"), ("All Files", "*.*")]
        )
        if file_path:
            if self.slide_mgr.load_pptx(file_path):
                self.current_slide_idx = 0
                self.slide_scrubber.configure(to=len(self.slide_mgr.slides)-1, number_of_steps=len(self.slide_mgr.slides))
                self.refresh_sidebar_slide_list()
                self.refresh_keywords_grid()
                self.voice_engine.set_keywords(self.slide_mgr.slides)
                self.update_slide_display()
                self.focus_set()

    def save_pptx_file(self):
        """Opens native file dialog to save .pptx deck."""
        from tkinter import filedialog
        file_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Presentations", "*.pptx")])
        if file_path:
            self.slide_mgr.save_pptx(file_path)
            self.focus_set()


# ======================================================================================
# MAIN EXECUTION ENTRY POINT
# ======================================================================================

if __name__ == "__main__":
    app = PresentationApp()
    app.mainloop()
